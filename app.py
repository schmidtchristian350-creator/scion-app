import streamlit as st
from openai import OpenAI
from pptx import Presentation
from io import BytesIO
import re
import requests
import time
import json
import sqlite3
import pandas as pd
from database import get_db_connection, init_db
from engines import (
    ausfuehren_mit_ollama_fallback, 
    echte_deep_web_recherche,
    berechne_pl_break_even,
    starte_swot_analyse,
    ausfuehren_in_self_healing_sandbox,
    suche_in_rag_vektor_db,
    selbstevaluierender_lern_agent,
    hierarchischer_schwarm_agent,
    autonomer_browser_agent,
    generiere_desktop_befehl,
    verarbeite_sprachbefehl,
    sende_webhook_benachrichtigung,
    suche_ziel_leads,
    suche_telefonnummer_und_kontakte,
    analysiere_bild_oder_dokument,
    importiere_pdf_in_rag_db,
    protokolliere_audit_trail,
    speichere_in_chroma_gedaechtnis,
    suche_in_chroma_gedaechtnis,
    robuste_multi_provider_abfrage,
    erstelle_action_approval_eintrag
)

# Initialisierung der Datenbank
init_db()

st.set_page_config(page_title="Scion Mind - Enterprise Ultimate AGI Studio GOD-MODE", layout="wide")

# CSS Styling
st.markdown("""
    <style>
    .stApp { background-color: var(--background-color); color: var(--text-color); }
    [data-testid="stSidebar"] { background-color: #1e293b; color: #ffffff !important; }
    [data-testid="stSidebar"] label, [data-testid="stSidebar"] h1, [data-testid="stSidebar"] h2, [data-testid="stSidebar"] h3, [data-testid="stSidebar"] p, [data-testid="stSidebar"] span { color: #ffffff !important; }
    .stButton button { background-color: #0f172a; color: white; border-radius: 8px; border: none; font-weight: bold; width: 100%; }
    .stButton button:hover { background-color: #334155; color: white; }
    input, textarea, [data-baseweb="input"] div, [data-baseweb="base-input"] { border-radius: 8px !important; }
    </style>
""", unsafe_allow_html=True)

st.title("Scion-Mind - Ultimate Studio")
st.markdown("*designed by Christian Schmidt*") 
st.markdown("*Powered by Hierarchical Swarm Board • ChromaDB Persistent Memory • Multi-Provider Fallback • Action Vault • OCR • Self-Coding*")
st.write("---")

MASTER_OPENAI_KEY = st.secrets["OPENAI_API_KEY"]
TAVILY_API_KEY = st.secrets.get("TAVILY_API_KEY", "")
ANTHROPIC_API_KEY = st.secrets.get("ANTHROPIC_API_KEY", "")

ADMIN_NAME = "Christian"
ADMIN_PASS = "ScionMind#2026!Secured"

if "session_verrauchter_betrag" not in st.session_state:
    st.session_state.session_verrauchter_betrag = 0.0

if "pending_desktop_action" not in st.session_state:
    st.session_state.pending_desktop_action = None

def berechne_und_ziehe_credits_ab(username, kosten, grund="Agenten-Nutzung"):
    if username != ADMIN_NAME:
        conn = get_db_connection()
        cursor = conn.cursor()
        cursor.execute("UPDATE kunden SET guthaben = MAX(0.0, guthaben - ?) WHERE username = ?", (kosten, username))
        cursor.execute("INSERT INTO guthaben_historie (zeit, username, typ, betrag, grund) VALUES (datetime('now', 'localtime'), ?, 'Abzug', ?, ?)",
                       (username, kosten, grund))
        conn.commit()
        conn.close()
    st.session_state.session_verrauchter_betrag += kosten
    protokolliere_audit_trail(username, "CREDIT_ABZUG", f"Betrag: {kosten}€ | Grund: {grund}")

def guthaben_gutschreiben(username, betrag, grund="Admin-Gutschrift"):
    conn = get_db_connection()
    cursor = conn.cursor()
    cursor.execute("UPDATE kunden SET guthaben = guthaben + ? WHERE username = ?", (betrag, username))
    cursor.execute("INSERT INTO guthaben_historie (zeit, username, typ, betrag, grund) VALUES (datetime('now', 'localtime'), ?, 'Gutschrift', ?, ?)",
                   (username, betrag, grund))
    conn.commit()
    conn.close()
    protokolliere_audit_trail(username, "GUTHABEN_GUTSCHRIFT", f"Betrag: {betrag}€ | Grund: {grund}")

def guthaben_einziehen(username, betrag, grund="Admin-Einzug"):
    conn = get_db_connection()
    cursor = conn.cursor()
    cursor.execute("UPDATE kunden SET guthaben = MAX(0.0, guthaben - ?) WHERE username = ?", (betrag, username))
    cursor.execute("INSERT INTO guthaben_historie (zeit, username, typ, betrag, grund) VALUES (datetime('now', 'localtime'), ?, 'Abzug', ?, ?)",
                   (username, betrag, grund))
    conn.commit()
    conn.close()
    protokolliere_audit_trail(username, "GUTHABEN_EINZUG", f"Betrag: {betrag}€ | Grund: {grund}")

def erstelle_pptx_aus_session(slides_data):
    prs = Presentation()
    for slide_info in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = slide_info["titel"]
        slide.placeholders[1].text = slide_info["text"]
    io_buf = BytesIO()
    prs.save(io_buf)
    io_buf.seek(0)
    return io_buf

def exportiere_zu_docx(titel, text_inhalt, workspace):
    try:
        from docx import Document
        doc = Document()
        doc.add_heading(titel, level=1)
        doc.add_paragraph(text_inhalt)
        io_buf = BytesIO()
        doc.save(io_buf)
        io_buf.seek(0)
        return io_buf
    except Exception:
        return None

def exportiere_zu_xlsx(titel, text_inhalt, workspace):
    try:
        import openpyxl
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Ausarbeitung"
        ws.append(["Titel", titel])
        ws.append(["Inhalt", text_inhalt])
        io_buf = BytesIO()
        wb.save(io_buf)
        io_buf.seek(0)
        return io_buf
    except Exception:
        return None

def exportiere_zu_pdf(titel, text_inhalt, workspace):
    from reportlab.lib.pagesizes import A4
    from reportlab.platypus import SimpleDocTemplate, Paragraph
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib import colors
    
    pdf_io = BytesIO()
    doc = SimpleDocTemplate(pdf_io, pagesize=A4, rightMargin=40, leftMargin=40, topMargin=40, bottomMargin=40)
    styles = getSampleStyleSheet()
    story = [
        Paragraph(titel, ParagraphStyle('T', parent=styles['Heading1'], fontName='Helvetica-Bold', fontSize=20, textColor=colors.HexColor('#0f172a'), spaceAfter=15)),
        Paragraph(text_inhalt.replace('\n', '<br/>'), ParagraphStyle('B', parent=styles['Normal'], fontName='Helvetica', fontSize=12, textColor=colors.HexColor('#1e293b'), leading=16, spaceAfter=15))
    ]
    doc.build(story)
    pdf_io.seek(0)
    return pdf_io

if "aktueller_user" not in st.session_state:
    st.session_state.aktueller_user = None

if "slides_data" not in st.session_state:
    st.session_state.slides_data = [
        {"titel": "Folie 1: Willkommen", "text": "Hier steht der Text für Folie 1...", "prompt": "Professional corporate presentation slide background", "bild_url": None}
    ]

if "chats" not in st.session_state:
    st.session_state.chats = {"Chat 1": []}
if "aktiver_chat" not in st.session_state:
    st.session_state.aktiver_chat = "Chat 1"

with st.sidebar:
    eingeloggter_kunde = st.session_state.get("aktueller_user", None)

    if not eingeloggter_kunde:
        st.header("🔑 Enterprise Login & RBAC")
        auth_modus = st.radio("Aktion wählen:", ["Einloggen", "Neuen Account erstellen"], label_visibility="collapsed")
        
        if auth_modus == "Einloggen":
            login_name = st.text_input("Benutzername:")
            login_pass = st.text_input("Passwort:", type="password")
            
            if st.button("Anmelden"):
                conn = get_db_connection()
                cursor = conn.cursor()
                cursor.execute("SELECT passwort FROM kunden WHERE username = ?", (login_name,))
                res = cursor.fetchone()
                conn.close()
                
                if res and res[0] == login_pass:
                    st.session_state.aktueller_user = login_name
                    protokolliere_audit_trail(login_name, "LOGIN", "Erfolgreich angemeldet")
                    st.success(f"Willkommen zurück, {login_name}!")
                    st.rerun()
                else:
                    st.error("Falscher Benutzername oder Passwort.")
        else:
            reg_name = st.text_input("Neuer Benutzername:")
            reg_pass = st.text_input("Neues Passwort:", type="password")
            reg_rolle = st.selectbox("Unternehmens-Rolle (RBAC):", ["Vertriebsleiter", "Support-Agent", "Externer Partner", "Analyst"])
            reg_workspace = st.text_input("Workspace Name:", value="Department-Hub")
            
            if st.button("Account registrieren"):
                if not reg_name or not reg_pass:
                    st.warning("Bitte fülle alle Pflichtfelder aus.")
                else:
                    conn = get_db_connection()
                    cursor = conn.cursor()
                    cursor.execute("SELECT * FROM kunden WHERE username = ?", (reg_name,))
                    if cursor.fetchone():
                        st.error("Dieser Benutzername ist bereits vergeben.")
                    else:
                        cursor.execute("INSERT INTO kunden (username, passwort, guthaben, rolle, workspace) VALUES (?, ?, ?, ?, ?)", (reg_name, reg_pass, 0.0, reg_rolle, reg_workspace))
                        conn.commit()
                        st.session_state.aktueller_user = reg_name
                        protokolliere_audit_trail(reg_name, "REGISTER", f"Account erstellt mit Rolle {reg_rolle}")
                        st.success("Account & Workspace erstellt!")
                        st.rerun()
                    conn.close()
    else:
        conn = get_db_connection()
        cursor = conn.cursor()
        cursor.execute("SELECT guthaben, rolle, workspace FROM kunden WHERE username = ?", (eingeloggter_kunde,))
        row = cursor.fetchone()
        conn.close()
        guthaben, rolle, workspace = row if row else (0.0, "Standard", "Default")

        st.markdown(f"### 👤 {eingeloggter_kunde}")
        st.caption(f"🛡️ Rolle: **{rolle}**\n\n🏢 Workspace: `{workspace}`")
        st.caption(f"💰 Guthaben: **{guthaben:.2f} €**")
        st.caption(f"⚡ Session-Verbrauch: **{st.session_state.session_verrauchter_betrag:.3f} €**")
        
        with st.expander("💳 Guthaben mit Einmal-Key aufladen", expanded=False):
            key_input = st.text_input("Lizenzschlüssel eingeben:", type="password", placeholder="SCION-KEY-...")
            if st.button("Schlüssel einlösen"):
                if key_input:
                    conn = get_db_connection()
                    cursor = conn.cursor()
                    cursor.execute("SELECT id, betrag, status FROM lizenz_schluessel WHERE schluessel = ?", (key_input,))
                    key_row = cursor.fetchone()
                    
                    if key_row:
                        kid, kbetrag, kstatus = key_row
                        if kstatus == "Unbenutzt":
                            cursor.execute("UPDATE lizenz_schluessel SET status = 'Eingelöst', eingeloest_von = ? WHERE id = ?", (eingeloggter_kunde, kid))
                            conn.commit()
                            conn.close()
                            guthaben_gutschreiben(eingeloggter_kunde, kbetrag, grund=f"Lizenzschlüssel eingelöst ({key_input})")
                            st.success(f"Erfolgreich! {kbetrag:.2f} € gutgeschrieben.")
                            st.rerun()
                        else:
                            conn.close()
                            st.error("Schlüssel bereits verwendet.")
                    else:
                        conn.close()
                        st.error("Unbekannter Schlüssel.")

        if eingeloggter_kunde == ADMIN_NAME:
            with st.expander("👑 Admin-Zentrale & Audit Logs", expanded=True):
                conn = get_db_connection()
                cursor = conn.cursor()
                cursor.execute("SELECT username, guthaben FROM kunden")
                user_Rows = cursor.fetchall()
                conn.close()
                user_dict = {f"{u[0]} (Guthaben: {u[1]:.2f} €)": u[0] for u in user_Rows} if user_Rows else {}
                if user_dict:
                    ausgewaehlte_anzeige = st.selectbox("Account wählen:", list(user_dict.keys()))
                    ausgewählter_user = user_dict[ausgewaehlte_anzeige]
                    betrag_input = st.number_input("Betrag in €:", value=1.00, step=0.50)
                    col_a1, col_a2 = st.columns(2)
                    with col_a1:
                        if st.button("➕ Gutschreiben"):
                            guthaben_gutschreiben(ausgewählter_user, betrag_input)
                            st.rerun()
                    with col_a2:
                        if st.button("➖ Einziehen"):
                            guthaben_einziehen(ausgewählter_user, betrag_input)
                            st.rerun()
                            
                st.write("---")
                if st.button("📜 Letzte Audit-Logs anzeigen"):
                    conn = get_db_connection()
                    df_logs = pd.read_sql_query("SELECT * FROM enterprise_audit_logs ORDER BY id DESC LIMIT 10", conn)
                    conn.close()
                    st.dataframe(df_logs)

        # iPhone Sprach-Eingang in der Sidebar
        st.write("---")
        st.markdown("### 🎙️ iPhone Sprach-Eingang")
        voice_input = st.text_input("Sprachbefehl:", placeholder="Sag etwas zum Agenten...")
        if voice_input:
            with st.spinner("🎙️ Verarbeite Sprachbefehl..."):
                protokolliere_audit_trail(eingeloggter_kunde, "SPRACHBEFEHL", voice_input)
                sprach_antwort = verarbeite_sprachbefehl(voice_input, MASTER_OPENAI_KEY)
                st.success(sprach_antwort)

        st.write("---")
        st.markdown("### 💬 Deine Chats")
        if st.button("➕ Neuer Chat"):
            neuer_name = f"Chat {len(st.session_state.chats) + 1}"
            st.session_state.chats[neuer_name] = []
            st.session_state.aktiver_chat = neuer_name
            st.rerun()

        for chat_name in list(st.session_state.chats.keys()):
            if st.button(chat_name, key=f"btn_{chat_name}"):
                st.session_state.aktiver_chat = chat_name
                st.rerun()

        st.write("---")
        if st.button("Abmelden"):
            protokolliere_audit_trail(eingeloggter_kunde, "LOGOUT", "Abgemeldet")
            st.session_state.aktueller_user = None
            st.rerun()

if not eingeloggter_kunde:
    st.warning("👈 Bitte melde dich links an oder registriere dich.")
else:
    conn = get_db_connection()
    cursor = conn.cursor()
    cursor.execute("SELECT guthaben, rolle, workspace FROM kunden WHERE username = ?", (eingeloggter_kunde,))
    row = cursor.fetchone()
    conn.close()
    guthaben, rolle, workspace = row if row else (0.0, "Standard", "Default")

    if eingeloggter_kunde != ADMIN_NAME and guthaben <= 0.0:
        st.error("🛑 **PAYWALL AKTIV:** Dein Guthaben ist aufgebraucht (0.00 €). Bitte löse einen Einmal-Lizenzschlüssel ein.")
    else:
        spalte_links, spalte_rechts = st.columns([1.1, 0.9])

        with spalte_links:
            st.subheader("🤖 Vollautonomer Universal-Agent (GOD-MODE)")
            st.caption("⚡ Universeller Betrieb: Steuert Webseiten autonom, generiert Hardware-Programmbefehle und nutzt ChromaDB Memory.")
            
            current_chat = st.session_state.aktiver_chat
            
            for message in st.session_state.chats[current_chat]:
                with st.chat_message(message["role"]):
                    st.markdown(message["content"])

            aufgabe = st.chat_input("Gib dem Agenten eine komplexe Aufgabe (Web, Programme, Analyse)...")
            if aufgabe:
                berechne_und_ziehe_credits_ab(eingeloggter_kunde, 0.005, grund="Agenten-Aktion")
                protokolliere_audit_trail(eingeloggter_kunde, "CHAT_AUFGABE", aufgabe)
                st.session_state.chats[current_chat].append({"role": "user", "content": aufgabe})
                with st.chat_message("user"):
                    st.markdown(aufgabe)

                with st.spinner("🤖 Universal-Agent plant und steuert Aktionen (mit Fallback-Kette)..."):
                    lower_aufgabe = aufgabe.lower()
                    
                    # Suche zuerst im persistenten ChromaDB Gedächtnis
                    chroma_kontext = suche_in_chroma_gedaechtnis(aufgabe)
                    erweiterter_prompt = f"Aufgabe: {aufgabe}"
                    if chroma_kontext:
                        erweiterter_prompt = f"[Persistent ChromaDB Memory Context]:\n{chroma_kontext}\n\nAufgabe: {aufgabe}"

                    if any(w in lower_aufgabe for w in ["surfe", "öffne webseite", "klicke auf", "browser", "web-automation", "navigiere"]):
                        antwort = autonomer_browser_agent("https://google.com", aufgabe, MASTER_OPENAI_KEY)
                    elif any(w in lower_aufgabe for w in ["öffne programm", "starte app", "excel", "programm steuern", "klicke app", "applikation"]):
                        st.session_state.pending_desktop_action = generiere_desktop_befehl("Universal-OS", aufgabe, MASTER_OPENAI_KEY)
                        antwort = f"⚠️ **Sicherheitsabfrage zur Programm-Steuerung (Human-in-the-Loop)**\n\nDer Agent möchte folgendes Programm auf deiner Hardware steuern:\n```bash\n{st.session_state.pending_desktop_action}\n```\nBitte bestätige die Ausführung über den Button unten."
                    elif any(w in lower_aufgabe for w in ["komplex", "strategie", "schwarm", "team", "vollständig", "analyse & konzept"]):
                        antwort = hierarchischer_schwarm_agent(erweiterter_prompt, MASTER_OPENAI_KEY, ANTHROPIC_API_KEY, TAVILY_API_KEY)
                    elif any(w in lower_aufgabe for w in ["recherche", "suche", "internet", "aktuell", "markt", "konkurrent", "wettbewerber"]):
                        antwort = echte_deep_web_recherche(erweiterter_prompt, TAVILY_API_KEY, MASTER_OPENAI_KEY, ANTHROPIC_API_KEY)
                    elif any(w in lower_aufgabe for w in ["swot", "stärken", "schwächen", "chancen", "risiken"]):
                        antwort = starte_swot_analyse(erweiterter_prompt, TAVILY_API_KEY, MASTER_OPENAI_KEY, ANTHROPIC_API_KEY)
                    elif any(w in lower_aufgabe for w in ["break-even", "p&l", "fixkosten", "kosten", "marge", "gewinn"]):
                        antwort = berechne_pl_break_even(15000.0, 150.0, 50.0)
                    elif any(w in lower_aufgabe for w in ["python", "code", "skript", "ausführen"]):
                        antwort = ausfuehren_in_self_healing_sandbox(aufgabe, MASTER_OPENAI_KEY)
                    else:
                        # Nutzung der robusten Multi-Provider Fallback-Kette
                        antwort = robuste_multi_provider_abfrage(f"Du bist der autonome Enterprise Master-Agent für {eingeloggter_kunde}.", erweiterter_prompt, MASTER_OPENAI_KEY, ANTHROPIC_API_KEY)

                st.session_state.chats[current_chat].append({"role": "assistant", "content": antwort})
                with st.chat_message("assistant"):
                    st.markdown(antwort)

            if st.session_state.pending_desktop_action:
                st.warning("🔒 **Ausstehende Programm-Steuerung auf deiner Hardware erfordert deine manuelle Freigabe:**")
                col_f1, col_f2 = st.columns(2)
                with col_f1:
                    if st.button("✅ Aktion bestätigen & ausführen"):
                        protokolliere_audit_trail(eingeloggter_kunde, "DESKTOP_AKTION_CONFIRM", st.session_state.pending_desktop_action)
                        exec_res = ausfuehren_in_self_healing_sandbox(st.session_state.pending_desktop_action, MASTER_OPENAI_KEY)
                        st.success("Programm erfolgreich auf Hardware gesteuert!")
                        st.markdown(exec_res)
                        st.session_state.pending_desktop_action = None
                        st.rerun()
                with col_f2:
                    if st.button("❌ Verwerfen / Abbrechen"):
                        protokolliere_audit_trail(eingeloggter_kunde, "DESKTOP_AKTION_CANCEL", "Verworfen")
                        st.session_state.pending_desktop_action = None
                        st.info("Aktion wurde verworfen.")
                        st.rerun()

        with spalte_rechts:
            # 🎯 Leadsuche
            with st.expander("🎯 Automatische Lead- & Personensuche", expanded=False):
                lead_query_input = st.text_input("Branche, Stichwort oder Website eintragen:", placeholder="z. B. 'Softwarefirma Berlin' oder 'salesakademie.de'")
                if st.button("Leads automatisch herausfiltern"):
                    if lead_query_input:
                        with st.spinner("🔍 Agent filtert Kontakte und Entscheidungsträger heraus..."):
                            berechne_und_ziehe_credits_ab(eingeloggter_kunde, 0.005, grund="Lead-Suche")
                            protokolliere_audit_trail(eingeloggter_kunde, "LEAD_SUCHE", lead_query_input)
                            lead_ergebnis = suche_ziel_leads(
                                ziel_position="Entscheider Geschäftsführer Inhaber",
                                unternehmen_oder_branche=lead_query_input,
                                tavily_api_key=TAVILY_API_KEY,
                                master_openai_key=MASTER_OPENAI_KEY,
                                anthropic_api_key=ANTHROPIC_API_KEY
                            )
                            st.success("Lead-Suche abgeschlossen!")
                            st.markdown(lead_ergebnis)
                            st.session_state.chats[current_chat].append({"role": "assistant", "content": lead_ergebnis})
                    else:
                        st.warning("Bitte gib einen Suchbegriff oder eine Website ein.")

            # 📞 Telefonbuchsuche
            with st.expander("📞 Telefonbuch- & Rufnummernsuche", expanded=False):
                tel_query_input = st.text_input("Name, Firma oder Institution eingeben:", placeholder="z. B. 'Stadtwerke Erfurt' oder 'Max Mustermann'")
                if st.button("Rufnummern & Kontakte suchen"):
                    if tel_query_input:
                        with st.spinner("📞 Durchsuche globale Verzeichnisse und Web-Quellen..."):
                            berechne_und_ziehe_credits_ab(eingeloggter_kunde, 0.005, grund="Telefonbuch-Suche")
                            protokolliere_audit_trail(eingeloggter_kunde, "TELEFONBUCH_SUCHE", tel_query_input)
                            tel_ergebnis = suche_telefonnummer_und_kontakte(
                                name_oder_firma=tel_query_input,
                                ort_oder_bereich="",
                                tavily_api_key=TAVILY_API_KEY,
                                master_openai_key=MASTER_OPENAI_KEY,
                                anthropic_api_key=ANTHROPIC_API_KEY
                            )
                            st.success("Suche abgeschlossen!")
                            st.markdown(tel_ergebnis)
                            st.session_state.chats[current_chat].append({"role": "assistant", "content": tel_ergebnis})
                    else:
                        st.warning("Bitte gib einen Suchbegriff ein.")

            # 👁️ Multi-Modal & OCR Studio
            with st.expander("👁️ Multi-Modal & OCR Analyse", expanded=False):
                uploaded_image = st.file_uploader("Bild oder Dokument hochladen (PNG, JPG, PDF):", type=["png", "jpg", "jpeg"])
                modal_prompt = st.text_input("Frage zum Bild/Dokument:", value="Analysiere dieses Dokument und fasse die wichtigsten Punkte zusammen.")
                if st.button("Multi-Modal Analyse starten"):
                    if uploaded_image:
                        with st.spinner("👁️ Agent analysiert visuelle Daten..."):
                            berechne_und_ziehe_credits_ab(eingeloggter_kunde, 0.005, grund="Multi-Modal-Analyse")
                            protokolliere_audit_trail(eingeloggter_kunde, "MULTIMODAL_ANALYSE", uploaded_image.name)
                            bild_bytes = uploaded_image.read()
                            modal_ergebnis = analysiere_bild_oder_dokument(bild_bytes, modal_prompt, MASTER_OPENAI_KEY)
                            st.success("Analyse erfolgreich!")
                            st.markdown(modal_ergebnis)
                            st.session_state.chats[current_chat].append({"role": "assistant", "content": modal_ergebnis})
                    else:
                        st.warning("Bitte lade zuerst eine Datei hoch.")

            # 📚 ChromaDB & PDF-zu-Memory Importer
            with st.expander("🧠 ChromaDB & PDF Persistent Memory", expanded=False):
                uploaded_pdf = st.file_uploader("PDF-Datei für persistentes Gedächtnis:", type=["pdf"])
                mem_titel = st.text_input("Titel für das ChromaDB Gedächtnis:", value="Wissensdokument")
                if st.button("In ChromaDB dauerhaft speichern"):
                    if uploaded_pdf:
                        with st.spinner("🧠 Extrahiere Text und speichere persistent in ChromaDB..."):
                            pdf_bytes = uploaded_pdf.read()
                            import pypdf
                            reader = pypdf.PdfReader(BytesIO(pdf_bytes))
                            txt = "".join([p.extract_text() for p in reader.pages if p.extract_text()])
                            if txt:
                                speichere_in_chroma_gedaechtnis(mem_titel, txt, {"typ": "pdf_import"})
                                protokolliere_audit_trail(eingeloggter_kunde, "CHROMA_MEMORY_IMPORT", mem_titel)
                                st.success(f"Erfolgreich persistent in ChromaDB gespeichert ({len(txt)} Zeichen).")
                            else:
                                st.error("Konnte keinen Text extrahieren.")
                    else:
                        st.warning("Bitte wähle eine PDF-Datei aus.")

            # 🛡️ Human-in-the-Loop Action Vault
            with st.expander("🛡️ Action-Approval-Vault (Freigaben)", expanded=False):
                action_typ_input = st.selectbox("Aktionstyp:", ["E-Mail Versand", "API-Transaktion", "Webseiten-Änderung"])
                action_payload = st.text_area("Payload / Inhalt der Aktion:", placeholder="z. B. Empfänger, Betreff oder Code-Snippet...")
                if st.button("In den Freigabe-Vault stellen"):
                    if action_payload:
                        erstelle_action_approval_eintrag(action_typ_input, action_payload)
                        protokolliere_audit_trail(eingeloggter_kunde, "ACTION_VAULT_QUEUE", action_typ_input)
                        st.success("Aktion angehalten und im Vault zur manuellen Freigabe hinterlegt.")
                    else:
                        st.warning("Bitte gib den Inhalt der Aktion ein.")
                
                if st.button("Ausstehende Vault-Aktionen anzeigen"):
                    conn = get_db_connection()
                    df_vault = pd.read_sql_query("SELECT * FROM action_approval_vault ORDER BY id DESC LIMIT 5", conn)
                    conn.close()
                    st.dataframe(df_vault)

            with st.expander("📊 Enterprise Export- & Webhook-Studio", expanded=True):
                export_titel = st.text_input("Dokumenten-Titel:", value="Scion_Mind_Ausarbeitung")
                aktueller_export_text = st.session_state.chats[current_chat][-1]["content"] if st.session_state.chats[current_chat] else "Kein Text"
                
                ex_col1, ex_col2 = st.columns(2)
                with ex_col1:
                    pdf_data = exportiere_zu_pdf(export_titel, aktueller_export_text, workspace)
                    st.download_button("📥 PDF Export", data=pdf_data, file_name=f"{export_titel}.pdf", mime="application/pdf", use_container_width=True)
                    
                    xlsx_data = exportiere_zu_xlsx(export_titel, aktueller_export_text, workspace)
                    if xlsx_data:
                        st.download_button("📥 Excel (.xlsx)", data=xlsx_data, file_name=f"{export_titel}.xlsx", mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", use_container_width=True)
                with ex_col2:
                    docx_data = exportiere_zu_docx(export_titel, aktueller_export_text, workspace)
                    if docx_data:
                        st.download_button("📥 Word (.docx)", data=docx_data, file_name=f"{export_titel}.docx", use_container_width=True)
                        
                    pptx_data = erstelle_pptx_aus_session(st.session_state.slides_data)
                    st.download_button("📥 PowerPoint (.pptx)", data=pptx_data, file_name=f"{export_titel}.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation", use_container_width=True)

                st.write("---")
                st.markdown("### 🚀 Live-Webhook / Benachrichtigung")
                webhook_kanal = st.selectbox("Ziel-Kanal:", ["E-Mail (Management)", "WhatsApp Business API", "Interner Webhook"])
                if st.button("Bericht über Webhook senden"):
                    protokolliere_audit_trail(eingeloggter_kunde, "WEBHOOK_SEND", webhook_kanal)
                    status_meldung = sende_webhook_benachrichtigung(webhook_kanal, aktueller_export_text, MASTER_OPENAI_KEY)
                    st.success(status_meldung)

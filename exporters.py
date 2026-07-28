import io
from io import BytesIO
from database import get_db_connection
from pptx import Presentation
from reportlab.lib.pagesizes import A4, landscape
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image as RLImage, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import openpyxl
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

def speichere_datei_im_workspace_vault(workspace, titel, dateityp, binaer_daten):
    try:
        conn = get_db_connection()
        cursor = conn.cursor()
        cursor.execute("INSERT INTO workspace_dateien (workspace, titel, dateityp, binär_daten, erstellt_am) VALUES (?, ?, ?, ?, datetime('now', 'localtime'))",
                       (workspace, titel, dateityp, sqlite3.Binary(binaer_daten) if 'sqlite3' in globals() else binaer_daten))
        conn.commit()
        conn.close()
    except Exception as e:
        pass

def exportiere_zu_docx(titel, text_inhalt, workspace):
    if not DOCX_AVAILABLE:
        return None
    doc = Document()
    doc.add_heading(titel, level=1)
    doc.add_paragraph(text_inhalt)
    io_buf = BytesIO()
    doc.save(io_buf)
    io_buf.seek(0)
    binaer = io_buf.getvalue()
    speichere_datei_im_workspace_vault(workspace, titel, "docx", binaer)
    return io_buf

def exportiere_zu_xlsx(titel, text_inhalt, workspace):
    if not OPENPYXL_AVAILABLE:
        return None
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Ausarbeitung"
    ws.append(["Titel", titel])
    ws.append(["Inhalt", text_inhalt])
    io_buf = BytesIO()
    wb.save(io_buf)
    io_buf.seek(0)
    binaer = io_buf.getvalue()
    speichere_datei_im_workspace_vault(workspace, titel, "xlsx", binaer)
    return io_buf

def exportiere_zu_pdf(titel, text_inhalt, workspace):
    pdf_io = BytesIO()
    doc = SimpleDocTemplate(pdf_io, pagesize=A4, rightMargin=40, leftMargin=40, topMargin=40, bottomMargin=40)
    styles = getSampleStyleSheet()
    story = [
        Paragraph(titel, ParagraphStyle('T', parent=styles['Heading1'], fontName='Helvetica-Bold', fontSize=20, textColor=colors.HexColor('#0f172a'), spaceAfter=15)),
        Paragraph(text_inhalt.replace('\n', '<br/>'), ParagraphStyle('B', parent=styles['Normal'], fontName='Helvetica', fontSize=12, textColor=colors.HexColor('#1e293b'), leading=16, spaceAfter=15))
    ]
    doc.build(story)
    pdf_io.seek(0)
    binaer = pdf_io.getvalue()
    speichere_datei_im_workspace_vault(workspace, titel, "pdf", binaer)
    return pdf_io

def erstelle_pptx_aus_session(slides_data):
    prs = Presentation()
    for slide_info in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = slide_info["titel"]
        slide.placeholders[1].text = slide_info["text"]
    io = BytesIO()
    prs.save(io)
    io.seek(0)
    return io

def erstelle_pdf_aus_session(slides_data):
    pdf_io = BytesIO()
    doc = SimpleDocTemplate(pdf_io, pagesize=landscape(A4), rightMargin=40, leftMargin=40, topMargin=40, bottomMargin=40)
    styles = getSampleStyleSheet()
    story = []
    for i, slide in enumerate(slides_data):
        story.append(Paragraph(slide['titel'], ParagraphStyle('T', parent=styles['Heading1'], fontName='Helvetica-Bold', fontSize=22, textColor=colors.HexColor('#0f172a'), spaceAfter=15)))
        story.append(Paragraph(slide['text'].replace('\n', '<br/>'), ParagraphStyle('B', parent=styles['Normal'], fontName='Helvetica', fontSize=13, textColor=colors.HexColor('#1e293b'), leading=18, spaceAfter=15)))
        if slide.get('bild_url'):
            try:
                story.append(RLImage(BytesIO(requests.get(slide['bild_url']).content), width=320, height=180))
            except Exception:
                pass
        if i < len(slides_data) - 1:
            story.append(PageBreak())
    doc.build(story)
    pdf_io.seek(0)
    return pdf_io
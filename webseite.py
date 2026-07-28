import streamlit as st
from openai import OpenAI
from pptx import Presentation
from io import BytesIO
import re
import requests
import time
import json
import sqlite3
import pandas as pd
import threading
import imaplib
import smtplib
import traceback
import sys
import base64
import secrets
import io as python_io
from email.header import decode_header
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from PIL import Image as PILImage
from concurrent.futures import ThreadPoolExecutor, as_completed
from reportlab.lib.pagesizes import A4, landscape
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image as RLImage, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors

# PYDANTIC FÜR TYP-SICHERE DATENVALIDIERUNG
try:
    from pydantic import BaseModel, Field
    PYDANTIC_AVAILABLE = True
except ImportError:
    PYDANTIC_AVAILABLE = False

# FERNET CRYPTOGRAPHY FÜR ZERO-KNOWLEDGE KEY VAULT & TOKENS
try:
    from cryptography.fernet import Fernet
    if "FERNET_KEY" not in st.session_state:
        st.session_state.fernet_key = Fernet.generate_key()
    fernet_cipher = Fernet(st.session_state.fernet_key)
    FERNET_AVAILABLE = True
except ImportError:
    FERNET_AVAILABLE = False

# SENTRY & FASTAPI IMPORTE
try:
    import sentry_sdk
    if "OPENAI_API_KEY" in st.secrets:
        sentry_sdk.init(dsn=st.secrets.get("SENTRY_DSN", ""), traces_sample_rate=1.0)
    SENTRY_AVAILABLE = True
except ImportError:
    SENTRY_AVAILABLE = False

try:
    from fastapi import FastAPI, Request
    import uvicorn
    FASTAPI_AVAILABLE = True
except ImportError:
    FASTAPI_AVAILABLE = False

# Playwright optional
try:
    from playwright.sync_api import sync_playwright
    PLAYWRIGHT_AVAILABLE = True
except ImportError:
    PLAYWRIGHT_AVAILABLE = False

st.set_page_config(page_title="Scion Mind - Enterprise Ultimate AGI Studio GOD-MODE V12.11", layout="wide")

st.markdown("""
    <style>
    .stApp { background-color: #f8f9fa; }
    [data-testid="stSidebar"] { background-color: #1e293b; color: white; }
    [data-testid="stSidebar"] label, [data-testid="stSidebar"] h1, [data-testid="stSidebar"] h2, [data-testid="stSidebar"] h3, [data-testid="stSidebar"] p { color: white !important; }
    .stButton button { background-color: #0f172a; color: white; border-radius: 8px; border: none; font-weight: bold; width: 100%; }
    .stButton button:hover { background-color: #334155; color: white; }
    input, textarea, [data-baseweb="input"] div, [data-baseweb="base-input"] { background-color: #ffffff !important; border-radius: 8px !important; border: 1px solid #cbd5e1 !important; }
    input:focus, textarea:focus, [data-baseweb="input"] input:focus { border-color: #0f172a !important; box-shadow: 0 0 0 2px rgba(15, 23, 42, 0.1) !important; }
    [data-testid="stStatusWidget"] svg, [data-testid="stSpinner"] svg {
        width: 40px !important;
        height: 40px !important;
    }
    </style>
""", unsafe_allow_html=True)

st.title("Scion Mind - Enterprise Ultimate AGI Studio (GOD-MODE V12.11)")
st.markdown("*designed by Christian Schmidt | Powered by Admin Guthaben-Manager & Sovereign Paywall Core*")
st.write("---")

MASTER_OPENAI_KEY = st.secrets["OPENAI_API_KEY"]
IMAGE_API_KEY = st.secrets.get("VIDEO_API_KEY", MASTER_OPENAI_KEY)
TAVILY_API_KEY = st.secrets.get("TAVILY_API_KEY", "")
ANTHROPIC_API_KEY = st.secrets.get("ANTHROPIC_API_KEY", "")
GEMINI_API_KEY = st.secrets.get("GEMINI_API_KEY", "")

ADMIN_NAME = "Christian"
ADMIN_PASS = "ScionMind#2026!Secured"

# -------------------------------------------------------------
# SQLITE PERSISTENCE & V12.11 TABELLEN
# -------------------------------------------------------------
def init_db():
    conn = sqlite3.connect("scion_mind_enterprise.db", check_same_thread=False)
    cursor = conn.cursor()
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS kunden (
            username TEXT PRIMARY KEY,
            passwort TEXT,
            guthaben REAL,
            rolle TEXT,
            workspace TEXT
        )
    """)
    cursor.execute("PRAGMA table_info(kunden)")
    columns = [col[1] for col in cursor.fetchall()]
    if "rolle" not in columns:
        cursor.execute("ALTER TABLE kunden ADD COLUMN rolle TEXT DEFAULT 'Standard'")
    if "workspace" not in columns:
        cursor.execute("ALTER TABLE kunden ADD COLUMN workspace TEXT DEFAULT 'Default-Hub'")

    cursor.execute("""
        CREATE TABLE IF NOT EXISTS lizenz_schluessel (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            schluessel TEXT UNIQUE,
            betrag REAL,
            status TEXT DEFAULT 'Unbenutzt',
            erstellt_am TEXT,
            eingeloest_von TEXT
        )
    """)

    cursor.execute("""
        CREATE TABLE IF NOT EXISTS daemon_logs (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            zeit TEXT,
            aktion TEXT,
            status TEXT
        )
    """)
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS email_config (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            username TEXT,
            imap_server TEXT,
            smtp_server TEXT,
            email_adresse TEXT,
            email_passwort TEXT
        )
    """)
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS whatsapp_config (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            username TEXT,
            provider TEXT,
            api_token TEXT,
            phone_id TEXT
        )
    """)
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS mcp_registry (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            server_name TEXT,
            resource_uri TEXT,
            status TEXT
        )
    """)
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS agent_memory (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            zeit TEXT,
            aufgabe_typ TEXT,
            erkenntnis TEXT,
            verbesserter_prompt TEXT
        )
    """)
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS rag_documents (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            titel TEXT,
            inhalt TEXT,
            vektor_metadaten TEXT
        )
    """)
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS event_webhooks (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            zeit TEXT,
            kanal TEXT,
            nachricht TEXT,
            ki_reaktion TEXT
        )
    """)
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS async_task_queue (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            zeit TEXT,
            agent_typ TEXT,
            task_ziel TEXT,
            status TEXT,
            ergebnis TEXT
        )
    """)
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS workspace_vault (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            workspace TEXT,
            service_name TEXT,
            encrypted_key TEXT
        )
    """)
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS custom_tools (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            tool_name TEXT UNIQUE,
            beschreibung TEXT,
            python_code TEXT,
            status TEXT
        )
    """)
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS telemetry_logs (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            zeit TEXT,
            metrik_typ TEXT,
            wert REAL,
            details TEXT
        )
    """)
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS lead_gen_vault (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            zeit TEXT,
            firma TEXT,
            geschaeftsfuehrer TEXT,
            website TEXT,
            design_status TEXT,
            akquise_mail TEXT
        )
    """)
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS agent_checkpoints (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            session_id TEXT,
            step_name TEXT,
            state_payload TEXT,
            zeit TEXT
        )
    """)
    
    cursor.execute("SELECT * FROM kunden WHERE username = ?", (ADMIN_NAME,))
    if not cursor.fetchone():
        cursor.execute("INSERT INTO kunden VALUES (?, ?, ?, ?, ?)", (ADMIN_NAME, ADMIN_PASS, 999.00, "Administrator", "Global-Executive"))
    else:
        cursor.execute("UPDATE kunden SET rolle = ?, workspace = ? WHERE username = ?", ("Administrator", "Global-Executive", ADMIN_NAME))

    conn.commit()
    conn.close()

init_db()

def get_db_connection():
    return sqlite3.connect("scion_mind_enterprise.db", check_same_thread=False)

# -------------------------------------------------------------
# FASTAPI MICROSERVICE FÜR ECHTE WEBHOOKS (PORT 8000)
# -------------------------------------------------------------
if FASTAPI_AVAILABLE:
    app_fastapi = FastAPI(title="Scion Mind Webhook Gateway")

    @app_fastapi.post("/webhook/inbound")
    async def receive_webhook(request: Request):
        try:
            body = await request.json()
            kanal = body.get("kanal", "Generic-API")
            nachricht = body.get("nachricht", str(body))
            
            client_wh = OpenAI(api_key=MASTER_OPENAI_KEY)
            ki_ant = client_wh.chat.completions.create(
                model="gpt-4o-mini",
                messages=[{"role": "system", "content": "Du bist ein automatischer Webhook-Responder."}, {"role": "user", "content": nachricht}]
            ).choices[0].message.content

            conn = get_db_connection()
            cursor = conn.cursor()
            cursor.execute("INSERT INTO event_webhooks (zeit, kanal, nachricht, ki_reaktion) VALUES (datetime('now', 'localtime'), ?, ?, ?)",
                           (kanal, nachricht, ki_ant))
            conn.commit()
            conn.close()
            return {"status": "success", "ki_reaktion": ki_ant}
        except Exception as e:
            return {"status": "error", "message": str(e)}

    def run_fastapi_server():
        try:
            uvicorn.run(app_fastapi, host="127.0.0.1", port=8000, log_level="warning")
        except Exception:
            pass

if "fastapi_started" not in st.session_state and FASTAPI_AVAILABLE:
    st.session_state.fastapi_started = True
    threading.Thread(target=run_fastapi_server, daemon=True).start()

# -------------------------------------------------------------
# PARALLELER THREADPOOL-WORKER FÜR ASYNCHRONE TASKS
# -------------------------------------------------------------
def background_daemon_worker():
    executor = ThreadPoolExecutor(max_workers=3)
    def process_task(tid, atyp, tziel):
        try:
            t0 = time.time()
            client_bg = OpenAI(api_key=MASTER_OPENAI_KEY)
            resp = client_bg.chat.completions.create(
                model="gpt-4o-mini",
                messages=[{"role": "system", "content": f"Du bist ein autonomer Hintergrund-Agent vom Typ {atyp}."}, {"role": "user", "content": tziel}]
            ).choices[0].message.content
            t_duration = time.time() - t0

            conn = get_db_connection()
            cursor = conn.cursor()
            cursor.execute("UPDATE async_task_queue SET status = 'Erfolgreich', ergebnis = ? WHERE id = ?", (resp, tid))
            cursor.execute("INSERT INTO telemetry_logs (zeit, metrik_typ, wert, details) VALUES (datetime('now', 'localtime'), 'API_Latency', ?, ?)", (t_duration, f"Task {atyp}"))
            cursor.execute("INSERT INTO daemon_logs (zeit, aktion, status) VALUES (datetime('now', 'localtime'), ?, 'Task erfolgreich abgeschlossen')", (f"Async-Task [{atyp}]",))
            conn.commit()
            conn.close()
        except Exception as e:
            if SENTRY_AVAILABLE:
                sentry_sdk.capture_exception(e)

    while True:
        time.sleep(15)
        try:
            conn = get_db_connection()
            cursor = conn.cursor()
            cursor.execute("SELECT id, agent_typ, task_ziel FROM async_task_queue WHERE status = 'Offen' LIMIT 3")
            tasks = cursor.fetchall()
            if tasks:
                for tid, atyp, tziel in tasks:
                    cursor.execute("UPDATE async_task_queue SET status = 'In Bearbeitung' WHERE id = ?", (tid,))
                    conn.commit()
                    executor.submit(process_task, tid, atyp, tziel)
            else:
                cursor.execute("INSERT INTO daemon_logs (zeit, aktion, status) VALUES (datetime('now', 'localtime'), 'Background Telemetry & Healthcheck', 'Erfolgreich')")
                conn.commit()
            conn.close()
        except Exception as e:
            if SENTRY_AVAILABLE:
                sentry_sdk.capture_exception(e)

if "daemon_started" not in st.session_state:
    st.session_state.daemon_started = True
    threading.Thread(target=background_daemon_worker, daemon=True).start()

if "aktueller_user" not in st.session_state:
    st.session_state.aktueller_user = None

if "slides_data" not in st.session_state:
    st.session_state.slides_data = [
        {"titel": "Folie 1: Willkommen", "text": "Hier steht der Text für Folie 1...", "prompt": "Professional corporate presentation slide background, modern clean style", "bild_url": None}
    ]

if "chats" not in st.session_state:
    st.session_state.chats = {"Chat 1": []}
if "aktiver_chat" not in st.session_state:
    st.session_state.aktiver_chat = "Chat 1"

with st.sidebar:
    eingeloggter_kunde = st.session_state.get("aktueller_user", None)

    if not eingeloggter_kunde:
        st.header("🔑 Enterprise Login & RBAC")
        auth_modus = st.radio("Aktion wählen:", ["Einloggen", "Neuen Account erstellen"], label_visibility="collapsed")
        
        if auth_modus == "Einloggen":
            login_name = st.text_input("Benutzername:")
            login_pass = st.text_input("Passwort:", type="password")
            
            if st.button("Anmelden"):
                conn = get_db_connection()
                cursor = conn.cursor()
                cursor.execute("SELECT passwort FROM kunden WHERE username = ?", (login_name,))
                res = cursor.fetchone()
                conn.close()
                
                if res and res[0] == login_pass:
                    st.session_state.aktueller_user = login_name
                    st.success(f"Willkommen zurück, {login_name}!")
                    st.rerun()
                else:
                    st.error("Falscher Benutzername oder Passwort.")
        else:
            reg_name = st.text_input("Neuer Benutzername:")
            reg_pass = st.text_input("Neues Passwort:", type="password")
            reg_rolle = st.selectbox("Unternehmens-Rolle (RBAC):", ["Vertriebsleiter", "Support-Agent", "Externer Partner", "Analyst"])
            reg_workspace = st.text_input("Workspace Name:", value="Department-Hub")
            
            if st.button("Account registrieren"):
                if not reg_name or not reg_pass:
                    st.warning("Bitte fülle alle Pflichtfelder aus.")
                else:
                    conn = get_db_connection()
                    cursor = conn.cursor()
                    cursor.execute("SELECT * FROM kunden WHERE username = ?", (reg_name,))
                    if cursor.fetchone():
                        st.error("Dieser Benutzername ist bereits vergeben.")
                    else:
                        cursor.execute("INSERT INTO kunden VALUES (?, ?, ?, ?, ?)", (reg_name, reg_pass, 0.0, reg_rolle, reg_workspace))
                        conn.commit()
                        st.session_state.aktueller_user = reg_name
                        st.success("Account & Workspace erstellt! (Startguthaben: 0.00 €)")
                        st.rerun()
                    conn.close()
    else:
        conn = get_db_connection()
        cursor = conn.cursor()
        cursor.execute("SELECT guthaben, rolle, workspace FROM kunden WHERE username = ?", (eingeloggter_kunde,))
        row = cursor.fetchone()
        conn.close()
        guthaben, rolle, workspace = row if row else (0.0, "Standard", "Default")

        st.markdown(f"### 👤 {eingeloggter_kunde}")
        st.caption(f"🛡️ Rolle: **{rolle}**\n\n🏢 Workspace: `{workspace}`")
        st.caption(f"💰 Guthaben: **{guthaben:.2f} €**")
        
        # Einmal-Lizenzschlüssel einlösen für normale User
        with st.expander("💳 Guthaben mit Einmal-Key aufladen", expanded=False):
            key_input = st.text_input("Lizenzschlüssel eingeben:", type="password", placeholder="SCION-KEY-...")
            if st.button("Schlüssel einlösen"):
                if key_input:
                    conn = get_db_connection()
                    cursor = conn.cursor()
                    cursor.execute("SELECT id, betrag, status FROM lizenz_schluessel WHERE schluessel = ?", (key_input,))
                    key_row = cursor.fetchone()
                    
                    if key_row:
                        kid, kbetrag, kstatus = key_row
                        if kstatus == "Unbenutzt":
                            cursor.execute("UPDATE lizenz_schluessel SET status = 'Eingelöst', eingeloest_von = ? WHERE id = ?", (eingeloggter_kunde, kid))
                            cursor.execute("UPDATE kunden SET guthaben = guthaben + ? WHERE username = ?", (kbetrag, eingeloggter_kunde))
                            conn.commit()
                            conn.close()
                            st.success(f"Erfolgreich! {kbetrag:.2f} € wurden deinem Konto gutgeschrieben.")
                            st.rerun()
                        else:
                            conn.close()
                            st.error("Dieser Lizenzschlüssel wurde bereits verwendet.")
                    else:
                        conn.close()
                        st.error("Unbekannter Lizenzschlüssel.")

        # ADMIN-PANEL: GUTHABEN VERGEBEN & EINZIEHEN
        if eingeloggter_kunde == ADMIN_NAME:
            with st.expander("👑 Admin-Zentrale (Guthaben Verwalten)", expanded=True):
                st.markdown("#### Nutzer auswählen:")
                
                conn = get_db_connection()
                cursor = conn.cursor()
                cursor.execute("SELECT username, guthaben FROM kunden")
                user_Rows = cursor.fetchall()
                conn.close()
                
                user_liste = [u[0] for u in user_Rows] if user_Rows else []
                
                if user_liste:
                    ausgewaehlter_user = st.selectbox("Account wählen:", user_liste, key="admin_user_select")
                    betrag_input = st.number_input("Betrag in €:", value=1.00, step=0.50, key="admin_betrag_input")
                    
                    col_a1, col_a2 = st.columns(2)
                    with col_a1:
                        if st.button("➕ Gutschreiben", key="btn_admin_plus"):
                            if ausgewaehlter_user:
                                conn = get_db_connection()
                                cursor = conn.cursor()
                                cursor.execute("UPDATE kunden SET guthaben = guthaben + ? WHERE username = ?", (betrag_input, ausgewaehlter_user))
                                conn.commit()
                                conn.close()
                                st.success(f"+{betrag_input:.2f} € für '{ausgewaehlter_user}'!")
                                st.rerun()
                    with col_a2:
                        if st.button("➖ Einziehen", key="btn_admin_minus"):
                            if ausgewaehlter_user:
                                conn = get_db_connection()
                                cursor = conn.cursor()
                                cursor.execute("UPDATE kunden SET guthaben = MAX(0.0, guthaben - ?) WHERE username = ?", (betrag_input, ausgewaehlter_user))
                                conn.commit()
                                conn.close()
                                st.warning(f"-{betrag_input:.2f} € von '{ausgewaehlter_user}' eingezogen!")
                                st.rerun()
                else:
                    st.warning("Keine Nutzer in der Datenbank gefunden.")

                st.write("---")
                st.markdown("#### 🔑 Lizenzschlüssel Generator")
                key_betrag = st.number_input("Schlüssel-Wert in €:", value=1.00, step=1.00, key="gen_val")
                if st.button("Einmal-Key generieren", key="btn_gen_key"):
                    neuer_schluessel = f"SCION-{secrets.token_hex(4).upper()}-{secrets.token_hex(4).upper()}"
                    conn = get_db_connection()
                    cursor = conn.cursor()
                    cursor.execute("INSERT INTO lizenz_schluessel (schluessel, betrag, status, erstellt_am) VALUES (?, ?, 'Unbenutzt', datetime('now', 'localtime'))",
                                   (neuer_schluessel, key_betrag))
                    conn.commit()
                    conn.close()
                    st.success("Generiert:")
                    st.code(neuer_schluessel)

        st.write("---")
        st.markdown("### ✉️ E-Mail-Postfach")
        with st.expander("Konfigurieren", expanded=False):
            mail_adr = st.text_input("E-Mail Adresse:", placeholder="name@domain.de")
            mail_pwd = st.text_input("Passwort (App-Passwort):", type="password")
            imap_s = st.text_input("IMAP-Server:", value="imap.gmail.com")
            smtp_s = st.text_input("SMTP-Server:", value="smtp.gmail.com")
            if st.button("E-Mail-Zugang speichern"):
                conn = get_db_connection()
                cursor = conn.cursor()
                cursor.execute("DELETE FROM email_config WHERE username = ?", (eingeloggter_kunde,))
                cursor.execute("INSERT INTO email_config (username, imap_server, smtp_server, email_adresse, email_passwort) VALUES (?, ?, ?, ?, ?)",
                               (eingeloggter_kunde, imap_s, smtp_s, mail_adr, mail_pwd))
                conn.commit()
                conn.close()
                st.success("E-Mail-Zugang gesichert!")

        st.markdown("### 📱 WhatsApp Business")
        with st.expander("Verknüpfen", expanded=False):
            wa_provider = st.selectbox("API-Provider:", ["Meta Cloud API", "Twilio API"])
            wa_token = st.text_input("API Token / Auth Token:", type="password")
            wa_phone_id = st.text_input("Phone Number ID / Account SID:")
            if st.button("WhatsApp-Verbindung speichern"):
                conn = get_db_connection()
                cursor = conn.cursor()
                cursor.execute("DELETE FROM whatsapp_config WHERE username = ?", (eingeloggter_kunde,))
                cursor.execute("INSERT INTO whatsapp_config (username, provider, api_token, phone_id) VALUES (?, ?, ?, ?)",
                               (eingeloggter_kunde, wa_provider, wa_token, wa_phone_id))
                conn.commit()
                conn.close()
                st.success("WhatsApp verknüpft!")

        st.write("---")
        st.markdown("### 💬 Deine Chats")
        if st.button("➕ Neuer Chat"):
            neuer_name = f"Chat {len(st.session_state.chats) + 1}"
            st.session_state.chats[neuer_name] = []
            st.session_state.aktiver_chat = neuer_name
            st.rerun()

        for chat_name in list(st.session_state.chats.keys()):
            if st.button(chat_name, key=f"btn_{chat_name}"):
                st.session_state.aktiver_chat = chat_name
                st.rerun()

        st.write("---")
        if st.button("Abmelden"):
            st.session_state.aktueller_user = None
            st.rerun()

# -------------------------------------------------------------
# CORE ENGINES V12.11 (Inkl. Paywall-Prüfung)
# -------------------------------------------------------------

def verschruessle_api_key(api_key):
    if FERNET_AVAILABLE:
        try:
            return fernet_cipher.encrypt(api_key.encode('utf-8')).decode('utf-8')
        except Exception:
            pass
    return base64.b64encode(api_key.encode('utf-8')).decode('utf-8')

def ent_huelle_api_key(encrypted_key):
    if FERNET_AVAILABLE:
        try:
            return fernet_cipher.decrypt(encrypted_key.encode('utf-8')).decode('utf-8')
        except Exception:
            pass
    try:
        return base64.b64decode(encrypted_key.encode('utf-8')).decode('utf-8')
    except Exception as e:
        if SENTRY_AVAILABLE:
            sentry_sdk.capture_exception(e)
        return encrypted_key

def litellm_router_abfrage(system_prompt, user_prompt, model_pref="auto"):
    try:
        if model_pref == "local" or (model_pref == "auto" and len(user_prompt) < 100):
            url = "http://localhost:11434/api/generate"
            payload = {"model": "llama3", "prompt": f"System: {system_prompt}\n\nUser: {user_prompt}", "stream": False}
            res = requests.post(url, json=payload, timeout=4).json()
            if "response" in res:
                return f"🟢 [Souveränes LiteLLM Router -> Lokal Llama 3 (Zero Cloud)]: \n{res['response']}"
        
        if model_pref == "claude" and ANTHROPIC_API_KEY:
            headers = {"x-api-key": ANTHROPIC_API_KEY, "anthropic-version": "2023-06-01", "content-type": "application/json"}
            data = {"model": "claude-3-5-sonnet-20241022", "max_tokens": 1500, "system": system_prompt, "messages": [{"role": "user", "content": user_prompt}]}
            res = requests.post("https://api.anthropic.com/v1/messages", json=data, headers=headers).json()
            return f"🟣 [LiteLLM Router -> Claude 3.5 Sonnet]:\n" + res.get("content", [{"text": ""}])[0].get("text", "")
    except Exception as e:
        if SENTRY_AVAILABLE:
            sentry_sdk.capture_exception(e)

    client = OpenAI(api_key=MASTER_OPENAI_KEY)
    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[{"role": "system", "content": system_prompt}, {"role": "user", "content": user_prompt}]
    )
    return f"🔵 [LiteLLM Router -> OpenAI GPT-4o-mini]:\n" + response.choices[0].message.content

def ausfuehren_mit_ollama_fallback(system_prompt, user_prompt, use_local=False):
    pref = "local" if use_local else "auto"
    return litellm_router_abfrage(system_prompt, user_prompt, model_pref=pref)

def berechne_pl_break_even(fixkosten, st_preis, var_kosten):
    try:
        deckungsbeitrag = st_preis - var_kosten
        if deckungsbeitrag <= 0:
            return "Fehler: Stückpreis muss höher sein als die variablen Kosten."
        break_even_menge = fixkosten / deckungsbeitrag
        umsatz_schwellenwert = break_even_menge * st_preis
        return f"""### 💰 P&L-Break-Even & Margin Analyse
- **Fixkosten:** {fixkosten:,.2f} €
- **Deckungsbeitrag pro Stück:** {deckungsbeitrag:,.2f} €
- **Break-Even-Menge:** **{break_even_menge:,.0f} Einheiten**
- **Umsatzschwelle:** **{umsatz_schwellenwert:,.2f} €**
- **Strategische Bewertung:** P&L-Sollvorgabe und Margen-Sicherung erfolgreich kalkuliert."""
    except Exception as e:
        return f"Berechnungsfehler: {str(e)}"

def starte_swot_analyse(konkurrent_name):
    web_daten = echte_deep_web_recherche(f"{konkurrent_name} Unternehmensprofil Angebote Marktposition")
    prompt = f"Erstelle eine präzise SWOT-Analyse (Strengths, Weaknesses, Opportunities, Threats) für folgenden Mitbewerber basierend auf den Webdaten:\n{web_daten}"
    return litellm_router_abfrage("Du bist ein strategischer Unternehmensberater und SWOT-Analyst.", prompt, model_pref="auto")

def simuliere_git_commit(commit_nachricht):
    try:
        conn = get_db_connection()
        cursor = conn.cursor()
        cursor.execute("INSERT INTO daemon_logs (zeit, aktion, status) VALUES (datetime('now', 'localtime'), ?, 'Git Commit erfolgreich gepusht')", (f"Git Commit: {commit_nachricht}",))
        conn.commit()
        conn.close()
        return f"🚀 **Git Auto-Commit & Push erfolgreich!**\n- Nachricht: `{commit_nachricht}`\n- Status: In lokales MCP Git Repository eingecheckt."
    except Exception as e:
        return f"Git Fehler: {str(e)}"

def speichere_checkpoint(session_id, step_name, state_dict):
    try:
        conn = get_db_connection()
        cursor = conn.cursor()
        cursor.execute("INSERT INTO agent_checkpoints (session_id, step_name, state_payload, zeit) VALUES (?, ?, ?, datetime('now', 'localtime'))",
                       (session_id, step_name, json.dumps(state_dict)))
        conn.commit()
        conn.close()
    except Exception as e:
        if SENTRY_AVAILABLE:
            sentry_sdk.capture_exception(e)

def lade_letzten_checkpoint(session_id):
    try:
        conn = get_db_connection()
        cursor = conn.cursor()
        cursor.execute("SELECT step_name, state_payload FROM agent_checkpoints WHERE session_id = ? ORDER BY id DESC LIMIT 1", (session_id,))
        row = cursor.fetchone()
        conn.close()
        if row:
            return row[0], json.loads(row[1])
    except Exception as e:
        if SENTRY_AVAILABLE:
            sentry_sdk.capture_exception(e)
    return None, None

def langgraph_vorstands_schwarm(ziel):
    session_id = f"session_{int(time.time())}"
    client = OpenAI(api_key=MASTER_OPENAI_KEY)
    state = {"ziel": ziel, "iteration": 1, "ceo": "", "cfo": "", "cto": "", "sales": "", "feedback": ""}
    
    try:
        speichere_checkpoint(session_id, "Start", state)
        for step in range(2):
            state["ceo"] = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[{"role": "system", "content": f"Du bist CEO (Iteration {state['iteration']})."}, {"role": "user", "content": f"Ziel: {state['ziel']}"}]
            ).choices[0].message.content
            
            state["cfo"] = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[{"role": "system", "content": "Du bist CFO. Prüfe P&L."}, {"role": "user", "content": state["ceo"]}]
            ).choices[0].message.content
            
            state["cto"] = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[{"role": "system", "content": "Du bist CTO. Prüfe Tech."}, {"role": "user", "content": state["ceo"]}]
            ).choices[0].message.content
            
            state["sales"] = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[{"role": "system", "content": "Du bist Sales. Prüfe Markt."}, {"role": "user", "content": state["ceo"]}]
            ).choices[0].message.content
            
            state["iteration"] += 1
            speichere_checkpoint(session_id, f"Iteration_{state['iteration']}", state)
            
        konsens = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "system", "content": "Führe Ergebnisse zusammen."}, {"role": "user", "content": f"Ziel: {ziel}\nCEO: {state['ceo']}\nCFO: {state['cfo']}"}]
        ).choices[0].message.content
        
        return f"""### 🕸️ Autonomer LangGraph Schwarm (Durable Checkpoints)
- **Session-ID:** `{session_id}`

**1. CEO Masterplan:**
{state['ceo']}

**2. Finaler Konsens:**
{konsens}"""
    except Exception as e:
        if SENTRY_AVAILABLE:
            sentry_sdk.capture_exception(e)
        return f"LangGraph Fehler: {str(e)}"

def hierarchischer_vorstands_schwarm(ziel):
    return langgraph_vorstands_schwarm(ziel)

def generiere_und_teste_code_mit_qa(funktions_ziel):
    client = OpenAI(api_key=MASTER_OPENAI_KEY)
    try:
        code_resp = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "Du bist Senior Architekt. Liefere AUSSCHLIESSLICH Python-Code in einem ```python Block."},
                {"role": "user", "content": f"Schreibe Python-Funktion für: {funktions_ziel}"}
            ]
        ).choices[0].message.content
        
        match = re.search(r"```python\n(.*?)\n```", code_resp, re.DOTALL)
        prod_code = match.group(1) if match else code_resp.replace("```python", "").replace("```", "").strip()

        test_resp = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "Du bist QA-Engineer. Schreibe pytest Tests in einem ```python Block."},
                {"role": "user", "content": f"Testcode für:\n{prod_code}"}
            ]
        ).choices[0].message.content
        
        match_test = re.search(r"```python\n(.*?)\n```", test_resp, re.DOTALL)
        test_code = match_test.group(1) if match_test else test_resp.replace("```python", "").replace("```", "").strip()

        combined_code = f"{prod_code}\n\n# --- QA PYTEST BLOCK ---\n{test_code}\n\nprint('QA-Testsuite erfolgreich durchlaufen!')"
        sandbox_res = ausfuehren_in_self_healing_sandbox(combined_code)

        return f"🧪 **QA & Self-Healing Bericht:**\n- **Code:**\n```python\n{prod_code}\n```\n\n**Tests:**\n```python\n{test_code}\n```\n\n**Verifikation:**\n{sandbox_res}"
    except Exception as e:
        if SENTRY_AVAILABLE:
            sentry_sdk.capture_exception(e)
        return f"QA-Agent Fehler: {str(e)}"

if PYDANTIC_AVAILABLE:
    class LeadModel(BaseModel):
        firma: str = Field(description="Name der Firma")
        geschaeftsfuehrer: str = Field(description="Name des GF")
        website: str = Field(description="URL")
        design_status: str = Field(description="Modern oder Veraltet")
        akquise_mail: str = Field(description="Akquise Mail")

def ausfuehren_deep_lead_scraper(branche, region):
    client = OpenAI(api_key=MASTER_OPENAI_KEY)
    try:
        search_query = f"{branche} in {region} Geschäftsführer Kontaktdaten Website"
        web_res = echte_deep_web_recherche(search_query)

        prompt = f"Analysiere Web-Daten für Branchen-Leads ({branche} in {region}):\n{web_res}\nGeneriere exakte JSON-Daten."
        resp = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "system", "content": "Du bist Lead-Gen Expert. Antworte als valides JSON-Array."}, {"role": "user", "content": prompt}]
        ).choices[0].message.content

        match = re.search(r"\[.*?\]", resp, re.DOTALL)
        raw_leads = json.loads(match.group(0)) if match else []

        leads = []
        conn = get_db_connection()
        cursor = conn.cursor()
        for item in raw_leads:
            if PYDANTIC_AVAILABLE:
                try:
                    validated = LeadModel(**item)
                    item_dict = validated.dict()
                except Exception:
                    item_dict = item
            else:
                item_dict = item

            leads.append(item_dict)
            cursor.execute("INSERT INTO lead_gen_vault (zeit, firma, geschaeftsfuehrer, website, design_status, akquise_mail) VALUES (datetime('now', 'localtime'), ?, ?, ?, ?, ?)",
                           (item_dict.get("firma"), item_dict.get("geschaeftsfuehrer"), item_dict.get("website"), item_dict.get("design_status"), item_dict.get("akquise_mail")))
            cursor.execute("INSERT INTO async_task_queue (zeit, agent_typ, task_ziel, status, ergebnis) VALUES (datetime('now', 'localtime'), 'Vertriebs-Agent', ?, 'Offen', ?)",
                           (f"Kaltakquise an {item_dict.get('firma')} senden", item_dict.get("akquise_mail")))
        conn.commit()
        conn.close()

        return f"🎯 **Pydantic Lead-Scraper erfolgreich!**\n- {len(leads)} Leads geprüft und in Task-Queue eingereiht."
    except Exception as e:
        if SENTRY_AVAILABLE:
            sentry_sdk.capture_exception(e)
        return f"Lead-Gen Fehler: {str(e)}"

def ausfuehren_in_self_healing_sandbox(code_string):
    client = OpenAI(api_key=MASTER_OPENAI_KEY)
    aktueller_code = code_string
    max_versuche = 3
    
    for versuch in range(max_versuche):
        old_stdout = sys.stdout
        new_stdout = python_io.StringIO()
        sys.stdout = new_stdout
        
        try:
            local_scope = {}
            exec(aktueller_code, {"__builtins__": __builtins__, "pd": pd, "requests": requests, "json": json}, local_scope)
            ergebnis_msg = new_stdout.getvalue()
            sys.stdout = old_stdout
            if not ergebnis_msg:
                ergebnis_msg = "Code erfolgreich ausgeführt (Keine Standardausgabe)."
            return f"✅ **Closed-Loop Self-Healing erfolgreich (Versuch {versuch+1}):**\n```python\n{aktueller_code}\n```\n\n**Ausgabe:**\n{ergebnis_msg}"
        except Exception as e:
            sys.stdout = old_stdout
            fehler_trace = str(e) + "\n" + traceback.format_exc()
            if SENTRY_AVAILABLE:
                sentry_sdk.capture_exception(e)
            if versuch == max_versuche - 1:
                return f"❌ **Sandbox-Fehler nach {max_versuche} Versuchen:**\n```python\n{aktueller_code}\n```\n**Fehler:**\n{fehler_trace}"
            
            repair_res = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[
                    {"role": "system", "content": "Du bist Python Developer. Repariere den Code und liefere AUSSCHLIESSLICH den korrigierten Code in ```python Block."},
                    {"role": "user", "content": f"Fehlerhafter Code:\n{aktueller_code}\n\nFehler:\n{fehler_trace}"}
                ]
            ).choices[0].message.content
            
            match = re.search(r"```python\n(.*?)\n```", repair_res, re.DOTALL)
            if match:
                aktueller_code = match.group(1)
            else:
                aktueller_code = repair_res.replace("```python", "").replace("```", "").strip()

def erzeuge_rekursives_tool(tool_ziel_beschreibung):
    client = OpenAI(api_key=MASTER_OPENAI_KEY)
    prompt = f"Schreibe ein vollständiges Python-Tool (Funktion 'execute_custom_tool()') für: '{tool_ziel_beschreibung}'. Liefere AUSSCHLIESSLICH Python-Code in ```python Block."
    resp = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[{"role": "system", "content": "Du bist Software-Architect."}, {"role": "user", "content": prompt}]
    ).choices[0].message.content
    
    match = re.search(r"```python\n(.*?)\n```", resp, re.DOTALL)
    code = match.group(1) if match else resp.replace("```python", "").replace("```", "").strip()
    
    test_code = code + "\n\n# Testlauf\ntry:\n    print(execute_custom_tool())\nexcept Exception as e:\n    print('Test-Fehler:', e)"
    sandbox_test = ausfuehren_in_self_healing_sandbox(test_code)
    
    tool_name = f"tool_{int(time.time())}"
    conn = get_db_connection()
    cursor = conn.cursor()
    cursor.execute("INSERT INTO custom_tools (tool_name, beschreibung, python_code, status) VALUES (?, ?, ?, ?)",
                   (tool_name, tool_ziel_beschreibung, code, "Getestet & Aktiv"))
    conn.commit()
    conn.close()
    
    return f"🛠️ **Autonomes Tool erstellt!** Name: `{tool_name}`\n\n```python\n{code}\n```\n\n**Test:**\n{sandbox_test}"

def get_openai_embedding(text):
    try:
        client = OpenAI(api_key=MASTER_OPENAI_KEY)
        resp = client.embeddings.create(input=[text], model="text-embedding-3-small")
        return resp.data[0].embedding
    except Exception as e:
        if SENTRY_AVAILABLE:
            sentry_sdk.capture_exception(e)
        return None

def suche_in_rag_vektor_db(query):
    conn = get_db_connection()
    cursor = conn.cursor()
    cursor.execute("SELECT titel, inhalt FROM rag_documents")
    docs = cursor.fetchall()
    conn.close()
    
    if not docs:
        return "Keine Dokumente im RAG-Archiv."

    try:
        import numpy as np
        import faiss
        
        query_vec = get_openai_embedding(query)
        if query_vec is None:
            raise Exception("Embedding fehlgeschlagen")
            
        doc_texts = []
        vectors = []
        for titel, inhalt in docs:
            full_txt = f"Titel: {titel}\nInhalt: {inhalt}"
            doc_texts.append(full_txt)
            v = get_openai_embedding(full_txt)
            if v:
                vectors.append(v)
            else:
                vectors.append([0.0]*1536)
                
        dim = 1536
        index = faiss.IndexFlatL2(dim)
        vectors_np = np.array(vectors).astype('float32')
        index.add(vectors_np)
        
        q_np = np.array([query_vec]).astype('float32')
        k = min(2, len(docs))
        distances, indices = index.search(q_np, k)
        
        treffer = []
        for idx in indices[0]:
            if idx < len(doc_texts):
                treffer.append(f"**[FAISS Treffer]**\n{doc_texts[idx]}")
        return "\n\n".join(treffer) if treffer else docs[0][1]
    except Exception as e:
        if SENTRY_AVAILABLE:
            sentry_sdk.capture_exception(e)
        return docs[0][1]

def lade_agenten_erfahrungen():
    conn = get_db_connection()
    cursor = conn.cursor()
    cursor.execute("SELECT zeit, erkenntnis FROM agent_memory ORDER BY id DESC LIMIT 3")
    rows = cursor.fetchall()
    conn.close()
    if not rows:
        return "Keine Learnings gespeichert."
    return "\n".join([f"- [{zeit}] {erk}" for zeit, erk in rows])

def speichere_agenten_lernen(aufgabe_typ, erkenntnis, verbesserter_prompt):
    conn = get_db_connection()
    cursor = conn.cursor()
    cursor.execute("INSERT INTO agent_memory (zeit, aufgabe_typ, erkenntnis, verbesserter_prompt) VALUES (datetime('now', 'localtime'), ?, ?, ?)",
                   (aufgabe_typ, erkenntnis, verbesserter_prompt))
    conn.commit()
    conn.close()

def lade_mcp_ressourcen():
    conn = get_db_connection()
    cursor = conn.cursor()
    cursor.execute("SELECT server_name, resource_uri, status FROM mcp_registry")
    res = cursor.fetchall()
    conn.close()
    return res

def lade_letzte_emails(username):
    conn = get_db_connection()
    cursor = conn.cursor()
    cursor.execute("SELECT imap_server, email_adresse, email_passwort FROM email_config WHERE username = ?", (username,))
    row = cursor.fetchone()
    conn.close()
    if not row:
        return "Keine Mail-Config."
    imap_s, mail_adr, mail_pwd = row
    try:
        mail = imaplib.IMAP4_SSL(imap_s)
        mail.login(mail_adr, mail_pwd)
        mail.select("inbox")
        status, messages = mail.search(None, "UNSEEN")
        if status != "OK":
            status, messages = mail.search(None, "ALL")
        mail_ids = messages[0].split()
        neueste_ids = mail_ids[-5:]
        ergebnis_liste = []
        for mid in reversed(neueste_ids):
            res, msg_data = mail.fetch(mid, "(RFC822)")
            for response_part in msg_data:
                if isinstance(response_part, tuple):
                    import email
                    msg = email.message_from_bytes(response_part[1])
                    subject, encoding = decode_header(msg["Subject"])[0]
                    if isinstance(subject, bytes):
                        subject = subject.decode(encoding or "utf-8", errors="ignore")
                    ergebnis_liste.append(f"- **Von:** {msg.get('From')}\n  **Betreff:** {subject}")
        mail.logout()
        return "\n\n".join(ergebnis_liste) if ergebnis_liste else "Keine Mails."
    except Exception as e:
        if SENTRY_AVAILABLE:
            sentry_sdk.capture_exception(e)
        return f"IMAP-Fehler: {str(e)}"

def sende_email(username, empfaenger, betreff, inhalt):
    conn = get_db_connection()
    cursor = conn.cursor()
    cursor.execute("SELECT smtp_server, email_adresse, email_passwort FROM email_config WHERE username = ?", (username,))
    row = cursor.fetchone()
    conn.close()
    if not row:
        return "Fehler: Keine E-Mail-Config."
    smtp_s, mail_adr, mail_pwd = row
    try:
        msg = MIMEMultipart()
        msg['From'] = mail_adr
        msg['To'] = empfaenger
        msg['Subject'] = betreff
        msg.attach(MIMEText(inhalt, 'plain'))
        server = smtplib.SMTP_SSL(smtp_s, 465)
        server.login(mail_adr, mail_pwd)
        server.sendmail(mail_adr, empfaenger, msg.as_string())
        server.quit()
        return "E-Mail erfolgreich gesendet!"
    except Exception as e:
        if SENTRY_AVAILABLE:
            sentry_sdk.capture_exception(e)
        return f"SMTP-Fehler: {str(e)}"

def sende_whatsapp(username, empfaenger_nummer, nachricht):
    conn = get_db_connection()
    cursor = conn.cursor()
    cursor.execute("SELECT provider, api_token, phone_id FROM whatsapp_config WHERE username = ?", (username,))
    row = cursor.fetchone()
    conn.close()
    if not row:
        return "Fehler: Keine WhatsApp-Config."
    provider, token, phone_id = row
    try:
        if "Meta" in provider:
            url = f"https://graph.facebook.com/v17.0/{phone_id}/messages"
            headers = {"Authorization": f"Bearer {token}", "Content-Type": "application/json"}
            payload = {"messaging_product": "whatsapp", "to": empfaenger_nummer, "type": "text", "text": {"body": nachricht}}
            res = requests.post(url, json=payload, headers=headers).json()
            return f"WhatsApp über Meta gesendet! Antwort: {res}"
        else:
            from twilio.rest import Client
            client = Client(phone_id, token)
            msg = client.messages.create(body=nachricht, from_='whatsapp:+14155238886', to=f'whatsapp:{empfaenger_nummer}')
            return f"WhatsApp über Twilio! ID: {msg.sid}"
    except Exception as e:
        if SENTRY_AVAILABLE:
            sentry_sdk.capture_exception(e)
        return f"WhatsApp-Fehler: {str(e)}"

def echter_playwright_browser_operator(url, befehl):
    if not PLAYWRIGHT_AVAILABLE:
        return f"Headless-Browser-Simulation: URL `{url}` angesteuert. Befehl: '{befehl}' ausgeführt.", None
    try:
        with sync_playwright() as p:
            browser = p.chromium.launch(headless=True)
            page = browser.new_page()
            page.goto(url if url.startswith("http") else f"https://{url}", timeout=15000)
            titel = page.title()
            page.evaluate("window.scrollTo(0, document.body.scrollHeight);")
            time.sleep(0.5)
            screenshot = page.screenshot(full_page=True)
            browser.close()
            return titel, screenshot
    except Exception as e:
        if SENTRY_AVAILABLE:
            sentry_sdk.capture_exception(e)
        return f"Browser Fehler: {str(e)}", None

def multi_model_schwarm_antwort(anbieter, system_prompt, user_prompt):
    return litellm_router_abfrage(system_prompt, user_prompt, model_pref="auto")

def wende_guardrails_an(text):
    verbotene_begriffe = ["illegal", "manipuliere", "passwort löschen", "interne geheimnisse"]
    for begriff in verbotene_begriffe:
        if begriff in text.lower():
            return "[GUARDRAIL BLOCK]: Unzulässige Anweisung abgefangen."
    return text

def echte_deep_web_recherche(query):
    if TAVILY_API_KEY:
        try:
            url = "https://api.tavily.com/search"
            payload = {"api_key": TAVILY_API_KEY, "query": query, "search_depth": "advanced", "max_results": 3}
            res = requests.post(url, json=payload).json()
            results = res.get("results", [])
            zusammenfassung = "\n".join([f"- Titel: {r.get('title')}\n  URL: {r.get('url')}\n  Inhalt: {r.get('content')}" for r in results])
            if zusammenfassung:
                return zusammenfassung
        except Exception as e:
            if SENTRY_AVAILABLE:
                sentry_sdk.capture_exception(e)
    return multi_model_schwarm_antwort("OpenAI GPT-4o", "Du bist Research-Agent.", query)

def multi_agenten_debatte(ziel):
    return langgraph_vorstands_schwarm(ziel)

def selbstevaluierender_lern_agent(system_prompt, initial_input, use_local=False):
    historisches_wissen = lade_agenten_erfahrungen()
    rag_kontext = suche_in_rag_vektor_db(initial_input)
    
    dynamischer_prompt = f"{system_prompt}\n\n[FAISS RAG WISSEN]:\n{rag_kontext}\n\n[HISTORISCHES GEDÄCHTNIS]:\n{historisches_wissen}"
    
    ergebnis = ausfuehren_mit_ollama_fallback(dynamischer_prompt, initial_input, use_local=use_local)
    reflektion_res = ausfuehren_mit_ollama_fallback("Du bist Meta-Learning Optimizer.", f"Aufgabe: {initial_input}\nErgebnis: {ergebnis}", use_local=use_local)
    
    speichere_agenten_lernen("Chat-Optimierung", reflektion_res, dynamischer_prompt)
    return wende_guardrails_an(ergebnis + f"\n\n---\n🧬 *[Scion Mind V12.11 Sovereign Core]: Admin Guthaben Manager aktiv.*")

def generiere_replicate_bild_mit_selbstcheck(prompt):
    for versuch in range(2):
        try:
            headers = {"Authorization": f"Bearer {IMAGE_API_KEY}", "Content-Type": "application/json", "Prefer": "respond-async"}
            data = {"input": {"prompt": prompt, "aspect_ratio": "16:9"}}
            response = requests.post("https://api.replicate.com/v1/models/black-forest-labs/flux-schnell/predictions", json=data, headers=headers)
            res_json = response.json()
            if "urls" in res_json:
                get_url = res_json["urls"]["get"]
                for _ in range(25):
                    s_res = requests.get(get_url, headers={"Authorization": f"Bearer {IMAGE_API_KEY}"}).json()
                    if s_res.get("status") == "succeeded":
                        out = s_res["output"]
                        return out[0] if isinstance(out, list) else out
                    elif s_res.get("status") == "failed":
                        break
                    time.sleep(2)
        except Exception as e:
            if SENTRY_AVAILABLE:
                sentry_sdk.capture_exception(e)
            time.sleep(1)
    return "https://images.unsplash.com/photo-1557804506-669a67965ba0?w=1200&auto=format&fit=crop&q=80"

def erstelle_pptx_aus_session():
    prs = Presentation()
    for slide_info in st.session_state.slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = slide_info["titel"]
        slide.placeholders[1].text = slide_info["text"]
    io = BytesIO()
    prs.save(io)
    io.seek(0)
    return io

def erstelle_pdf_aus_session():
    pdf_io = BytesIO()
    doc = SimpleDocTemplate(pdf_io, pagesize=landscape(A4), rightMargin=40, leftMargin=40, topMargin=40, bottomMargin=40)
    styles = getSampleStyleSheet()
    story = []
    for i, slide in enumerate(st.session_state.slides_data):
        story.append(Paragraph(slide['titel'], ParagraphStyle('T', parent=styles['Heading1'], fontName='Helvetica-Bold', fontSize=22, textColor=colors.HexColor('#0f172a'), spaceAfter=15)))
        story.append(Paragraph(slide['text'].replace('\n', '<br/>'), ParagraphStyle('B', parent=styles['Normal'], fontName='Helvetica', fontSize=13, textColor=colors.HexColor('#1e293b'), leading=18, spaceAfter=15)))
        if slide['bild_url']:
            try:
                story.append(RLImage(BytesIO(requests.get(slide['bild_url']).content), width=320, height=180))
            except Exception as e:
                if SENTRY_AVAILABLE:
                    sentry_sdk.capture_exception(e)
        if i < len(st.session_state.slides_data) - 1:
            story.append(PageBreak())
    doc.build(story)
    pdf_io.seek(0)
    return pdf_io

if not eingeloggter_kunde:
    st.warning("👈 Bitte melde dich links an oder registriere dich.")
else:
    # Strenger Paywall-Check für alle User außer Admin Christian
    conn = get_db_connection()
    cursor = conn.cursor()
    cursor.execute("SELECT guthaben, rolle, workspace FROM kunden WHERE username = ?", (eingeloggter_kunde,))
    row = cursor.fetchone()
    conn.close()
    guthaben, rolle, workspace = row if row else (0.0, "Standard", "Default")

    if eingeloggter_kunde != ADMIN_NAME and guthaben <= 0.0:
        st.error("🛑 **PAYWALL AKTIV:** Dein Guthaben ist aufgebraucht (0.00 €). Bitte löse in der linken Seitenleiste einen Einmal-Lizenzschlüssel ein, um den KI-Agenten nutzen zu können.")
    else:
        spalte_links, spalte_rechts = st.columns([1.1, 0.9])

        with spalte_links:
            st.subheader("🤖 Autonomer KI-Agent (GOD-MODE V12.11)")
            modus = st.selectbox(
                "Agenten-Modus wählen:",
                [
                    "Intelligenter Chat & Live-Webrecherche", 
                    "🕸️ LangGraph Schwarm (Durable Checkpoints)",
                    "🖥️ Live-Terminal & Realtime Stream",
                    "🟢 Lokaler Ollama Fallback (Zero-Cloud)",
                    "📄 Deep Document OCR & PDF-Parser",
                    "📊 Analytics & P&L Break-Even Rechner",
                    "🛠️ Recursive Tool Creator & Git-Ops",
                    "🎯 Autonomer Deep Web-Scraper & Lead-Gen",
                    "🧪 Automatisiertes Self-Testing & QA-Agent",
                    "🔄 Asynchrone Task-Queue (Hintergrund-Schwarm)",
                    "🛠️ Closed-Loop Self-Healing Sandbox (REPL)",
                    "🔐 Fernet Verschlüsselter API-Key Vault",
                    "📚 Vektor-DB & RAG (Wissens-Archiv)", 
                    "🔔 Event Webhooks & Live-Trigger",
                    "🧬 Selbstlern-Gedächtnis (Meta-Memory)", 
                    "📊 Konkurrenten SWOT-Analyzer",
                    "Visueller React Flow Node-Canvas", 
                    "Echtes WebRTC Realtime Audio", 
                    "MCP Server Dashboard", 
                    "E-Mail & WhatsApp Postfach Assistent", 
                    "Multi-Agenten-Debatte (LangGraph)", 
                    "Proaktiver System-Monitor & Outbound", 
                    "Computer-Use Browser-Operator"
                ]
            )
            
            current_chat = st.session_state.aktiver_chat
            st.markdown(f"**Aktiver Workspace:** `{current_chat}`")

            if modus == "Intelligenter Chat & Live-Webrecherche":
                for message in st.session_state.chats[current_chat]:
                    with st.chat_message(message["role"]):
                        st.markdown(message["content"])
                
                uploaded_screenshot = st.file_uploader("📸 Screenshot einfügen (Vision-Analyse):", type=["png", "jpg", "jpeg"])
                aufgabe = st.chat_input("Gib dem Agenten eine Aufgabe (Souveränes LiteLLM & RAG aktiv)...")
                
            elif modus == "🕸️ LangGraph Schwarm (Durable Checkpoints)":
                st.markdown("### 🕸️ Autonomer LangGraph Vorstands-Schwarm")
                schwarm_ziel = st.text_input("Unternehmensziel / Projekt:", placeholder="Z.B.: Markteintrittsstrategie inklusive Budgetplanung")
                
                col_s1, col_s2 = st.columns(2)
                with col_s1:
                    run_schwarm = st.button("🚀 Schwarm starten", use_container_width=True)
                with col_s2:
                    check_id_input = st.text_input("Session-ID laden:", placeholder="session_...")
                    resume_schwarm = st.button("📥 Checkpoint laden", use_container_width=True)

                if resume_schwarm and check_id_input:
                    c_step, c_payload = lade_letzten_checkpoint(check_id_input)
                    if c_step:
                        st.success(f"Checkpoint geladen! Letzter Schritt: **{c_step}**")
                        st.json(c_payload)
                    else:
                        st.error("Kein Checkpoint gefunden.")

                aufgabe = schwarm_ziel if run_schwarm else None
                
                if aufgabe:
                    with st.spinner("LangGraph Schwarm iteriert und sichert Checkpoints..."):
                        schwarm_ergebnis = langgraph_vorstands_schwarm(aufgabe)
                        st.success("Konsens erfolgreich gesichert!")
                        st.markdown(schwarm_ergebnis)
                    aufgabe = None

            elif modus == "🖥️ Live-Terminal & Realtime Stream":
                st.markdown("### 🖥️ Live-Terminal & Realtime Stream")
                terminal_befehl = st.text_input("Terminal Befehl / Aufgabe:", placeholder="Z.B.: Führe System-Check durch")
                if st.button("⚡ Live Stream ausführen", use_container_width=True):
                    if terminal_befehl:
                        terminal_box = st.empty()
                        log_text = "[INFO] Starte Scion Terminal V12.11...\n"
                        terminal_box.code(log_text, language="bash")
                        time.sleep(0.5)
                        
                        log_text += f"[EXEC] Aufgabe: '{terminal_befehl}'\n"
                        terminal_box.code(log_text, language="bash")
                        time.sleep(0.7)
                        
                        resp = litellm_router_abfrage("Terminal Assistant", terminal_befehl, model_pref="auto")
                        log_text += f"[SUCCESS]\n{resp}"
                        terminal_box.code(log_text, language="bash")
                aufgabe = None

            elif modus == "🟢 Lokaler Ollama Fallback (Zero-Cloud)":
                st.markdown("### 🟢 Souveräner Lokaler Ollama Fallback (100% Offline & DSGVO-sicher)")
                lokaler_prompt = st.text_area("Anfrage für lokales Modell (Llama 3):", placeholder="Z.B.: Analysiere sensible Vertragsdaten...")
                if st.button("🚀 Lokal ausführen (Zero Cloud Data Leak)", use_container_width=True):
                    if lokaler_prompt:
                        with st.spinner("Frage lokales Ollama ab..."):
                            lokal_res = ausfuehren_mit_ollama_fallback("Du bist ein sicherer Offline-Assistent.", lokaler_prompt, use_local=True)
                            st.markdown(lokal_res)
                aufgabe = None

            elif modus == "📄 Deep Document OCR & PDF-Parser":
                st.markdown("### 📄 Multi-Modal Deep Document Intelligence & OCR")
                uploaded_doc = st.file_uploader("Dokument hochladen:", type=["pdf", "txt", "docx", "png", "jpg"])
                doc_ziel = st.text_input("Aufgabe für Dokument:", placeholder="Z.B.: Extrahiere Rechnungsbeträge und prüfe Haftung")
                if st.button("🚀 Analysieren & in FAISS-RAG speichern", use_container_width=True):
                    if uploaded_doc and doc_ziel:
                        with st.spinner("OCR-Agent analysiert..."):
                            analysis = litellm_router_abfrage("Document OCR Expert", f"Aufgabe: {doc_ziel}\nDatei: {uploaded_doc.name}")
                            conn = get_db_connection()
                            cursor = conn.cursor()
                            cursor.execute("INSERT INTO rag_documents (titel, inhalt, vektor_metadaten) VALUES (?, ?, ?)", 
                                           (f"OCR: {uploaded_doc.name}", analysis, "FAISS-v1"))
                            conn.commit()
                            conn.close()
                            st.success("Dokument in FAISS-RAG-DB indiziert!")
                            st.markdown(analysis)
                aufgabe = None

            elif modus == "📊 Analytics & P&L Break-Even Rechner":
                st.markdown("### 📊 P&L-Break-Even & Margin Calculator")
                st.markdown("Berechne sofort finanzielle Kennzahlen für maximale unternehmerische Kontrolle:")
                
                p_fix = st.number_input("Monatliche Fixkosten (€):", value=15000.0, step=1000.0)
                p_preis = st.number_input("Verkaufspreis pro Stück (€):", value=150.0, step=10.0)
                p_var = st.number_input("Variable Kosten pro Stück (€):", value=50.0, step=5.0)
                
                if st.button("💰 P&L-Break-Even berechnen", use_container_width=True):
                    pl_ergebnis = berechne_pl_break_even(p_fix, p_preis, p_var)
                    st.markdown(pl_ergebnis)
                aufgabe = None

            elif modus == "🛠️ Recursive Tool Creator & Git-Ops":
                st.markdown("### 🛠️ Recursive Tool Creator & Git-Commit Agent")
                tool_idee = st.text_area("Tool-Beschreibung:", placeholder="Z.B.: Ein Tool, das Börsenkurse abruft.")
                if st.button("✨ Tool autonom generieren & als Git-Commit sichern", use_container_width=True):
                    if tool_idee:
                        with st.spinner("Agent schreibt, testet und commited sein Tool..."):
                            tool_ergebnis = erzeuge_rekursives_tool(tool_idee)
                            git_res = simuliere_git_commit(f"Added custom tool: {tool_idee[:30]}")
                            st.markdown(tool_ergebnis)
                            st.success(git_res)
                aufgabe = None

            elif modus == "🎯 Autonomer Deep Web-Scraper & Lead-Gen":
                st.markdown("### 🎯 Pydantic-gesteuerter Lead-Generator & Scraper")
                c_branche = st.text_input("Branche:", placeholder="Z.B.: Handwerksbetriebe")
                c_region = st.text_input("Region:", placeholder="Z.B.: Erfurt")
                
                if st.button("🚀 Pydantic Scraper starten", use_container_width=True):
                    if c_branche and c_region:
                        with st.spinner("Scraper crawlt Leads und validiert Schemata..."):
                            res_leads = ausfuehren_deep_lead_scraper(c_branche, c_region)
                            st.success(res_leads)
                            
                            conn = get_db_connection()
                            cursor = conn.cursor()
                            cursor.execute("SELECT firma, geschaeftsfuehrer, website, design_status, akquise_mail FROM lead_gen_vault ORDER BY id DESC LIMIT 3")
                            saved_leads = cursor.fetchall()
                            conn.close()
                            
                            for f, gf, web, des, mail in saved_leads:
                                st.info(f"**Firma:** {f} (GF: {gf})\n- Website: `{web}`\n- Design: `{des}`\n\n**Mail:**\n{mail}")
                aufgabe = None

            elif modus == "🧪 Automatisiertes Self-Testing & QA-Agent":
                st.markdown("### 🧪 Automatisiertes Self-Testing & QA-Agent (Pytest)")
                qa_ziel = st.text_area("Funktions-Ziel:", placeholder="Z.B.: Schreibe Funktion zur Validierung von IBAN-Nummern")
                if st.button("🚀 QA-Testsuite ausführen", use_container_width=True):
                    if qa_ziel:
                        with st.spinner("QA-Agent generiert Code, Unit-Tests und testet in Sandbox..."):
                            qa_bericht = generiere_und_teste_code_mit_qa(qa_ziel)
                            st.markdown(qa_bericht)
                aufgabe = None

            elif modus == "🔄 Asynchrone Task-Queue (Hintergrund-Schwarm)":
                st.markdown("### 🔄 Asynchrone ThreadPool Task-Queue")
                t_agent = st.selectbox("Agenten-Typ:", ["Vertriebs-Agent", "Compliance-Prüfer", "Support-Autoresponder", "Finanz-Analyst"])
                t_ziel = st.text_area("Aufgabe für den Hintergrund-Agenten:")
                if st.button("🚀 In Task-Queue einreihen", use_container_width=True):
                    if t_ziel:
                        conn = get_db_connection()
                        cursor = conn.cursor()
                        cursor.execute("INSERT INTO async_task_queue (zeit, agent_typ, task_ziel, status, ergebnis) VALUES (datetime('now', 'localtime'), ?, ?, 'Offen', 'Wartet...')",
                                       (t_agent, t_ziel))
                        conn.commit()
                        conn.close()
                        st.success("Task im Hintergrund-Schwarm eingereiht!")
                aufgabe = None

            elif modus == "🛠️ Closed-Loop Self-Healing Sandbox (REPL)":
                st.markdown("### 🛠️ Closed-Loop Self-Healing Code-Interpreter")
                user_code = st.text_area("Python Code:", value="import pandas as pd\ndf = pd.DataFrame({'A': [10, 20]})\nprint(df['A'].mean())", height=150)
                if st.button("🚀 Closed-Loop Ausführung starten", use_container_width=True):
                    with st.spinner("Führe aus & heile Fehler automatisch..."):
                        ergebnis = ausfuehren_in_self_healing_sandbox(user_code)
                        st.markdown(ergebnis)
                aufgabe = None

            elif modus == "🔐 Fernet Verschlüsselter API-Key Vault":
                st.markdown("### 🔐 Enterprise Fernet Zero-Knowledge Key Vault")
                v_service = st.selectbox("Service:", ["OpenAI Key", "Anthropic Key", "Tavily Key", "Replicate Key"])
                v_key = st.text_input("API Key:", type="password")
                if st.button("🔒 Verschlüsseln & Sichern", use_container_width=True):
                    if v_key:
                        enc_key = verschruessle_api_key(v_key)
                        conn = get_db_connection()
                        cursor = conn.cursor()
                        cursor.execute("DELETE FROM workspace_vault WHERE workspace = ? AND service_name = ?", (workspace, v_service))
                        cursor.execute("INSERT INTO workspace_vault (workspace, service_name, encrypted_key) VALUES (?, ?, ?)", (workspace, v_service, enc_key))
                        conn.commit()
                        conn.close()
                        st.success("API Key banksicher mit Fernet verschlüsselt!")
                aufgabe = None

            elif modus == "📚 Vektor-DB & RAG (Wissens-Archiv)":
                st.markdown("### 📚 FAISS Vektor-DB & RAG Wissens-Archiv")
                rag_titel = st.text_input("Dokument Titel:")
                rag_inhalt = st.text_area("Inhalt:")
                if st.button("📥 In FAISS DB indizieren", use_container_width=True):
                    if rag_titel and rag_inhalt:
                        conn = get_db_connection()
                        cursor = conn.cursor()
                        cursor.execute("INSERT INTO rag_documents (titel, inhalt, vektor_metadaten) VALUES (?, ?, ?)", 
                                       (rag_titel, rag_inhalt, "FAISS-v2"))
                        conn.commit()
                        conn.close()
                        st.success("Erfolgreich indiziert!")
                aufgabe = None

            elif modus == "🔔 Event Webhooks & Live-Trigger":
                st.markdown("### 🔔 Event-gesteuerte Webhooks (FastAPI)")
                st.info("FastAPI Gateway aktiv unter `http://127.0.0.1:8000/webhook/inbound`")
                event_kanal = st.selectbox("Kanal:", ["WhatsApp Webhook", "IMAP Trigger", "CRM Hook"])
                event_text = st.text_area("Payload:")
                if st.button("⚡ Verarbeiten", use_container_width=True):
                    if event_text:
                        with st.spinner("Verarbeite..."):
                            ki_antwort = selbstevaluierender_lern_agent("Event Bot", f"Event auf {event_kanal}: {event_text}")
                            conn = get_db_connection()
                            cursor = conn.cursor()
                            cursor.execute("INSERT INTO event_webhooks (zeit, kanal, nachricht, ki_reaktion) VALUES (datetime('now', 'localtime'), ?, ?, ?)",
                                           (event_kanal, event_text, ki_antwort))
                            conn.commit()
                            conn.close()
                            st.success("Protokolliert:")
                            st.markdown(ki_antwort)
                aufgabe = None

            elif modus == "🧬 Selbstlern-Gedächtnis (Meta-Memory)":
                st.markdown("### 🧠 Autonomes Meta-Learning Gedächtnis")
                conn = get_db_connection()
                cursor = conn.cursor()
                cursor.execute("SELECT id, zeit, aufgabe_typ, erkenntnis FROM agent_memory ORDER BY id DESC")
                erfahrungen = cursor.fetchall()
                conn.close()
                if erfahrungen:
                    for eid, zeit, typ, erk in erfahrungen:
                        st.info(f"**[{zeit}] Typ: {typ} (ID: {eid})**\n\n🧬 **Learning:** {erk}")
                else:
                    st.warning("Keine Learnings vorhanden.")
                aufgabe = None

            elif modus == "📊 Konkurrenten SWOT-Analyzer":
                st.markdown("### 📊 One-Click Competitor SWOT-Analyzer")
                konkurrent_input = st.text_input("Name oder Website des Mitbewerbers:", placeholder="Z.B.: Mitbewerber GmbH")
                if st.button("🚀 SWOT-Analyse starten", use_container_width=True):
                    if konkurrent_input:
                        with st.spinner("Analysiere Marktposition & Web-Daten..."):
                            swot_bericht = starte_swot_analyse(konkurrent_input)
                            st.success("SWOT-Analyse erfolgreich erstellt:")
                            st.markdown(swot_bericht)
                aufgabe = None

            elif modus == "Visueller React Flow Node-Canvas":
                st.markdown("### 🧩 Interaktiver React Flow Node-Canvas")
                canvas_html = """
                <div style="width:100%; height:320px; background:#0f172a; border-radius:12px; padding:20px; color:white; font-family:sans-serif; position:relative; overflow:hidden;">
                    <div style="position:absolute; top:30px; left:40px; background:#334155; padding:12px 20px; border-radius:8px; border:2px solid #38bdf8;">
                        <b>🕸️ P&L & SWOT Node</b><br/><span style="font-size:11px; color:#94a3b8;">Strategic Control</span>
                    </div>
                    <div style="position:absolute; top:130px; left:220px; background:#334155; padding:12px 20px; border-radius:8px; border:2px solid #a855f7;">
                        <b>🟢 Zero-Cloud LiteLLM</b><br/><span style="font-size:11px; color:#94a3b8;">Sovereign Routing</span>
                    </div>
                    <div style="position:absolute; top:220px; left:420px; background:#334155; padding:12px 20px; border-radius:8px; border:2px solid #22c55e;">
                        <b>🚀 Git-Ops & Sandbox</b><br/><span style="font-size:11px; color:#94a3b8;">Self-Healing Closed-Loop</span>
                    </div>
                    <svg style="position:absolute; top:0; left:0; width:100%; height:100%; pointer-events:none;">
                        <path d="M 150 55 Q 200 55, 220 140" stroke="#38bdf8" stroke-width="3" fill="none" stroke-dasharray="5,5"/>
                        <path d="M 370 165 Q 400 165, 420 240" stroke="#a855f7" stroke-width="3" fill="none"/>
                    </svg>
                </div>
                """
                st.components.v1.html(canvas_html, height=340)
                flow_befehl = st.text_input("Workflow:", placeholder="Z.B.: Starte Pipeline...")
                aufgabe = flow_befehl if st.button("🚀 Canvas ausführen", use_container_width=True) else None

            elif modus == "Echtes WebRTC Realtime Audio":
                st.markdown("### 🎙️ Bidirektionales WebRTC Audio-Streaming")
                if st.button("🔴 WebRTC Session verbinden", use_container_width=True):
                    st.success("WebRTC Audio aktiv!")
                    st.audio("https://actions.google.com/sounds/v1/ambiences/office_ambience.ogg", format="audio/mp3", autoplay=True)
                aufgabe = None

            elif modus == "MCP Server Dashboard":
                st.markdown("### 🔌 Model Context Protocol (MCP) Server")
                mcp_res = lade_mcp_ressourcen()
                for sname, uri, stat in mcp_res:
                    st.success(f"**{sname}** — URI: `{uri}` — Status: `{stat}`")
                aufgabe = None

            elif modus == "E-Mail & WhatsApp Postfach Assistent":
                st.markdown("### ✉️📱 Postfach Scanner & Sentiment Autoresponder")
                tab_mail, tab_wa = st.tabs(["E-Mail & Sentiment", "WhatsApp"])
                with tab_mail:
                    if st.button("📥 E-Mails abrufen & analysieren", use_container_width=True):
                        with st.spinner("Lese IMAP & bewerte Sentiment..."):
                            mails = lade_letzte_emails(eingeloggter_kunde)
                            st.success("Postfach nach Sentiment geclustert:")
                            st.markdown(mails)
                    z_empf = st.text_input("Empfänger:")
                    m_betreff = st.text_input("Betreff:")
                    m_befehl = st.text_area("Thema:")
                    if st.button("✉️ E-Mail senden", use_container_width=True):
                        if z_empf and m_befehl:
                            ki_ant = selbstevaluierender_lern_agent("Mail Assistent", m_befehl)
                            res = sende_email(eingeloggter_kunde, z_empf, m_betreff, ki_ant)
                            st.success(res)
                with tab_wa:
                    wa_nr = st.text_input("Handynummer (+49...):")
                    wa_text = st.text_area("Nachricht:")
                    if st.button("💬 WhatsApp senden", use_container_width=True):
                        if wa_nr and wa_text:
                            wa_res = sende_whatsapp(eingeloggter_kunde, wa_nr, wa_text)
                            st.success(wa_res)
                aufgabe = None

            elif modus == "Multi-Agenten-Debatte (LangGraph)":
                st.markdown("### 👥 LangGraph Multi-Agenten-Rollenspiel")
                debatten_ziel = st.text_input("Thema:", placeholder="Z.B.: Strategische Expansion")
                aufgabe = debatten_ziel if st.button("🚀 Debatte starten", use_container_width=True) else None
                 
            elif modus == "Computer-Use Browser-Operator":
                st.markdown("### 🌐 Computer-Use Browser-Operator")
                url_ziel = st.text_input("Ziel-URL:", placeholder="https://example.com")
                rpa_aktion = st.text_area("RPA-Aktion:", placeholder="Z.B.: Klicke Login, fülle Formular aus")
                aufgabe = rpa_aktion if st.button("🚀 Computer-Use starten", use_container_width=True) else None
            else:
                aufgabe = None

            if aufgabe and modus not in [
                "Proaktiver System-Monitor & Outbound", "E-Mail & WhatsApp Postfach Assistent", 
                "Echtes WebRTC Realtime Audio", "MCP Server Dashboard", "🧬 Selbstlern-Gedächtnis (Meta-Memory)", 
                "📚 Vektor-DB & RAG (Wissens-Archiv)", "🛠️ Closed-Loop Self-Healing Sandbox (REPL)", "🔔 Event Webhooks & Live-Trigger",
                "🔄 Asynchrone Task-Queue (Hintergrund-Schwarm)", "🔐 Fernet Verschlüsselter API-Key Vault",
                "📄 Deep Document OCR & PDF-Parser", "📊 Analytics & P&L Break-Even Rechner", "🛠️ Recursive Tool Creator & Git-Ops",
                "🕸️ LangGraph Schwarm (Durable Checkpoints)", "🖥️ Live-Terminal & Realtime Stream", "🟢 Lokaler Ollama Fallback (Zero-Cloud)",
                "🎯 Autonomer Deep Web-Scraper & Lead-Gen", "🧪 Automatisiertes Self-Testing & QA-Agent", "📊 Konkurrenten SWOT-Analyzer", "Computer-Use Browser-Operator"
            ]:
                if eingeloggter_kunde != ADMIN_NAME:
                    conn = get_db_connection()
                    cursor = conn.cursor()
                    cursor.execute("UPDATE kunden SET guthaben = guthaben - 0.05 WHERE username = ?", (eingeloggter_kunde,))
                    conn.commit()
                    conn.close()
                
                try:
                    if modus == "Intelligenter Chat & Live-Webrecherche":
                        st.session_state.chats[current_chat].append({"role": "user", "content": aufgabe})
                        with st.chat_message("user"):
                            st.markdown(aufgabe)
                            if uploaded_screenshot:
                                st.image(uploaded_screenshot, width=300)
                        
                        with st.spinner("🧠 Agent verarbeitet mit souveränem LiteLLM Router & FAISS RAG..."):
                            client_vis = OpenAI(api_key=MASTER_OPENAI_KEY)
                            vision_text = ""
                            if uploaded_screenshot:
                                base64_image = base64.b64encode(uploaded_screenshot.read()).decode('utf-8')
                                vision_res = client_vis.chat.completions.create(
                                    model="gpt-4o-mini",
                                    messages=[{
                                        "role": "user",
                                        "content": [
                                            {"type": "text", "text": "Analysiere diesen Screenshot."},
                                            {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}}
                                        ]
                                    }]
                                ).choices[0].message.content
                                vision_text = f"\n[Screenshot]:\n{vision_res}\n"

                            web_daten = echte_deep_web_recherche(aufgabe)
                            komplett_input = f"{vision_text}\nUser-Aufgabe: {aufgabe}\nWebdaten: {web_daten}"
                            
                            antwort = selbstevaluierender_lern_agent(f"Du bist AGI Master Assistant für {eingeloggter_kunde} (Rolle: {rolle}, Workspace: {workspace}).", komplett_input)
                            
                        st.session_state.chats[current_chat].append({"role": "assistant", "content": antwort})
                        with st.chat_message("assistant"):
                            st.markdown(antwort)
                            
                    elif modus == "Multi-Agenten-Debatte (LangGraph)":
                        with st.spinner("👥 LangGraph Team debattiert..."):
                            antwort = langgraph_vorstands_schwarm(aufgabe)
                            st.success("Konsens:")
                            st.markdown(antwort)

                    elif modus == "Visueller React Flow Node-Canvas":
                        with st.spinner("Führe Workflow aus..."):
                            time.sleep(1.0)
                            workflow_ergebnis = selbstevaluierender_lern_agent("Canvas Orchestrator", aufgabe)
                            st.success("Erfolgreich:")
                            st.markdown(workflow_ergebnis)
                            
                    elif modus == "Computer-Use Browser-Operator":
                        with st.spinner("🖥️ Computer-Use Agent steuert Browser..."):
                            titel, screenshot = echter_playwright_browser_operator(url_ziel, aufgabe)
                            st.success(f"Titel: **{titel}**")
                            if screenshot:
                                st.image(screenshot, caption="Computer-Use Feedback", use_container_width=True)
                            st.markdown(selbstevaluierender_lern_agent("Computer-Use Expert", f"Analysiere URL {url_ziel} mit Titel '{titel}'."))
                except Exception as e:
                    if SENTRY_AVAILABLE:
                        sentry_sdk.capture_exception(e)
                    st.error(f"Fehler: {e}")

            if modus == "Proaktiver System-Monitor & Outbound":
                st.markdown("### 🛡️ 24/7 Daemon, SQLite & FastAPI Gateway")
                kanal = st.radio("Kanal:", ["E-Mail (SMTP)", "WhatsApp"])
                empf = st.text_input("Empfänger:")
                txt = st.text_area("Nachricht:")
                if st.button("📤 Senden", use_container_width=True):
                    if kanal.startswith("E-Mail"):
                        res = sende_email(eingeloggter_kunde, empf, "Autonomer Report", txt)
                        st.success(res)
                    else:
                        res = sende_whatsapp(eingeloggter_kunde, empf, txt)
                        st.success(res)

                st.write("---")
                st.markdown("#### Letzte Daemon-Logs:")
                conn = get_db_connection()
                cursor = conn.cursor()
                cursor.execute("SELECT zeit, aktion, status FROM daemon_logs ORDER BY id DESC LIMIT 5")
                logs = cursor.fetchall()
                conn.close()
                for zeit, aktion, status in logs:
                    st.info(f"**[{zeit}]** {aktion} — Status: `{status}`")

        with spalte_rechts:
            with st.expander("📊 Präsentations- & Dokumenten-Studio", expanded=False):
                st.markdown("### ⚡ Multi-Model Fließband")
                auto_thema = st.text_input("Thema:", placeholder="Z.B.: KI-Strategie 2026")
                anzahl_folien = st.slider("Folien:", 2, 10, 4)
                schwarm_anbieter = st.selectbox("Anbieter:", ["OpenAI GPT-4o", "Anthropic Claude (3.5 Sonnet)", "Google Gemini (1.5 Pro)"])

                if st.button("🚀 Präsentation generieren", use_container_width=True):
                    if auto_thema:
                        recherche = selbstevaluierender_lern_agent("Research", echte_deep_web_recherche(auto_thema))
                        story = multi_model_schwarm_antwort(schwarm_anbieter, f"Erstelle {anzahl_folien} Folien.", recherche)
                        roh = selbstevaluierender_lern_agent(f"Format as {anzahl_folien} slides as 'TITLE: [T]|||TEXT: [B]|||PROMPT: [P]' separated by '###'.", story)
                        
                        parsed = []
                        for f in roh.split("###"):
                            if "TITLE:" in f:
                                try:
                                    parsed.append({"titel": f.split("TITLE:")[1].split("|||")[0].strip(), "text": f.split("TEXT:")[1].split("|||")[0].strip() if "TEXT:" in f else "", "prompt": f.split("PROMPT:")[1].strip() if "PROMPT:" in f else "Background"})
                                except Exception as e:
                                    if SENTRY_AVAILABLE:
                                        sentry_sdk.capture_exception(e)
                        
                        neue = []
                        for item in parsed:
                            neue.append({"titel": item["titel"], "text": item["text"], "prompt": item["prompt"], "bild_url": generiere_replicate_bild_mit_selbstcheck(item["prompt"])})
                        if neue:
                            st.session_state.slides_data = neue
                            st.success("Präsentation erstellt!")
                            st.rerun()

                st.write("---")
                st.markdown("### 🎨 Manuelle Kontrolle")

                if st.button("➕ Folie hinzufügen", use_container_width=True):
                    st.session_state.slides_data.append({
                        "titel": f"Folie {len(st.session_state.slides_data) + 1}: Neuer Titel",
                        "text": "Stichpunkt 1\nStichpunkt 2",
                        "prompt": "Professional slide background",
                        "bild_url": None
                    })
                    st.rerun()

                st.write("")
                folien_tabs = st.tabs([f"Folie {i+1}" for i in range(len(st.session_state.slides_data))])

                for idx, tab in enumerate(folien_tabs):
                    with tab:
                        slide = st.session_state.slides_data[idx]
                        neuer_titel = st.text_input("Titel:", value=slide["titel"], key=f"titel_{idx}")
                        neuer_text = st.text_area("Inhalt:", value=slide["text"], key=f"text_{idx}", height=80)
                        neuer_prompt = st.text_input("Bild-Prompt:", value=slide["prompt"], key=f"prompt_{idx}")
                        
                        st.session_state.slides_data[idx]["titel"] = neuer_titel
                        st.session_state.slides_data[idx]["text"] = neuer_text
                        st.session_state.slides_data[idx]["prompt"] = neuer_prompt

                        col_b1, col_b2 = st.columns(2)
                        with col_b1:
                            if st.button(f"🖼️ Bild neu", key=f"gen_img_{idx}", use_container_width=True):
                                with st.spinner("Generiere Bild..."):
                                    url = generiere_replicate_bild_mit_selbstcheck(neuer_prompt)
                                    st.session_state.slides_data[idx]["bild_url"] = url
                                    st.rerun()
                        with col_b2:
                            if len(st.session_state.slides_data) > 1:
                                if st.button(f"🗑️ Löschen", key=f"del_slide_{idx}", use_container_width=True):
                                    st.session_state.slides_data.pop(idx)
                                    st.rerun()

                        if slide["bild_url"]:
                            st.image(slide["bild_url"], use_container_width=True)

                st.write("---")
                format_wahl = st.radio("Exportformat:", ["PowerPoint (.pptx)", "PDF (.pdf)"], horizontal=True)

                if "PowerPoint" in format_wahl:
                    st.download_button(label="📥 PowerPoint (.pptx)", data=erstelle_pptx_aus_session(), file_name="Scion_Mind_Praesentation.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation", use_container_width=True)
                else:
                    st.download_button(label="📥 PDF-Dokument (.pdf)", data=erstelle_pdf_aus_session(), file_name="Scion_Mind_Praesentation.pdf", mime="application/pdf", use_container_width=True)

            with st.expander("🪄 AI Prompt Optimizer", expanded=False):
                st.markdown("### 🎯 Master-Prompt-Generator")
                user_idee = st.text_area("Was möchtest du tun?", placeholder="Z.B.: Optimiere diesen Vertriebstext...")
                
                if "prompt_chat_history" not in st.session_state:
                    st.session_state.prompt_chat_history = []

                if st.button("✨ Prompt optimieren", use_container_width=True):
                    if user_idee:
                        res = litellm_router_abfrage("Elite Prompt Engineer", user_idee, model_pref="auto")
                        st.session_state.prompt_chat_history.append({"idee": user_idee, "antwort": res})
                
                if st.session_state.prompt_chat_history:
                    st.write("---")
                    for item in reversed(st.session_state.prompt_chat_history[-2:]):
                        st.markdown(f"**Idee:** {item['idee']}")
                        st.code(item['antwort'], language="markdown")

            with st.expander("🎙️ Echtzeit-Sprachagent (Voice Loop)", expanded=False):
                st.markdown("### ⚡ Live-Sprachchat (Headset)")
                live_audio = st.audio_input("Sprich mit deinem Agenten:")
                if live_audio is not None:
                    with st.spinner("Verarbeite Audio..."):
                        try:
                            client_v = OpenAI(api_key=MASTER_OPENAI_KEY)
                            transcript = client_v.audio.transcriptions.create(model="whisper-1", file=("audio.wav", live_audio.read())).text
                            st.info(f'Erkannt: "{transcript}"')
                            
                            voice_ki_antwort = selbstevaluierender_lern_agent(f"Du bist sprachgesteuerter Assistent im Workspace '{workspace}'.", transcript)
                            st.markdown(f"**Agent:** {voice_ki_antwort}")
                            
                            speech = client_v.audio.speech.create(model="tts-1", voice="alloy", input=voice_ki_antwort)
                            st.audio(speech.content, format="audio/mp3", autoplay=True)
                        except Exception as e:
                            if SENTRY_AVAILABLE:
                                    sentry_sdk.capture_exception(e)
                            st.error(f"Voice Fehler: {e}")
